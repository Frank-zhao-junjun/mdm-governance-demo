#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""AI 数据治理案例与趋势 PPT 生成器 — 20页商务风，白底深蓝，微软雅黑，含备注"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ---------- 主题 ----------
NAVY   = RGBColor(0x1F, 0x38, 0x64)   # 主色 深蓝
BLUE   = RGBColor(0x2E, 0x86, 0xAB)   # 强调 青蓝
ORANGE = RGBColor(0xE6, 0x7E, 0x22)   # 强调 橙
GREEN  = RGBColor(0x27, 0xAE, 0x60)   # 成功绿
RED    = RGBColor(0xC0, 0x39, 0x2B)   # 警示红
GRAY   = RGBColor(0x4A, 0x4A, 0x4A)   # 正文灰
LGRAY  = RGBColor(0x8A, 0x8A, 0x8A)   # 浅灰
BG     = RGBColor(0xF2, 0xF4, 0xF7)   # 卡片底
BORDER = RGBColor(0xD5, 0xDB, 0xE3)   # 卡片边
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
FONT   = "微软雅黑"

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]

def set_font(run, size=14, bold=False, color=GRAY, name=FONT):
    f = run.font
    f.size, f.bold, f.name = Pt(size), bold, name
    f.color.rgb = color
    # east asian font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', name)

def add_text(slide, x, y, w, h, text, size=14, bold=False, color=GRAY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.12):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        set_font(r, size, bold, color)
    return tb

def add_rect(slide, x, y, w, h, fill=BG, line=BORDER, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line_w=0.75):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp

def add_card(slide, x, y, w, h, title, body, accent=BLUE, tsize=14, bsize=11.5):
    add_rect(slide, x, y, w, h)
    add_rect(slide, x, y, Inches(0.07), h, fill=accent, line=None, shape=MSO_SHAPE.RECTANGLE)
    add_text(slide, x + Inches(0.18), y + Inches(0.10), w - Inches(0.3), Inches(0.4),
             title, size=tsize, bold=True, color=NAVY)
    add_text(slide, x + Inches(0.18), y + Inches(0.5), w - Inches(0.3), h - Inches(0.6),
             body, size=bsize, color=GRAY)

def add_slide_header(slide, section, title, page_no):
    add_rect(slide, 0, 0, SW, Inches(0.12), fill=NAVY, line=None, shape=MSO_SHAPE.RECTANGLE)
    add_rect(slide, Inches(0.55), Inches(0.42), Inches(0.09), Inches(0.55), fill=BLUE, line=None, shape=MSO_SHAPE.RECTANGLE)
    add_text(slide, Inches(0.78), Inches(0.35), Inches(9.5), Inches(0.3),
             section, size=12, bold=True, color=BLUE)
    add_text(slide, Inches(0.78), Inches(0.62), Inches(11.5), Inches(0.55),
             title, size=23, bold=True, color=NAVY)
    add_text(slide, Inches(12.35), Inches(0.40), Inches(0.8), Inches(0.3),
             str(page_no), size=13, bold=True, color=LGRAY, align=PP_ALIGN.RIGHT)

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def new_slide():
    return prs.slides.add_slide(BLANK)

def add_table(slide, x, y, w, h, data, col_w=None, header_fill=NAVY,
              fsize=11, header_size=11.5, row_h=0.32):
    rows, cols = len(data), len(data[0])
    gf = slide.shapes.add_table(rows, cols, x, y, w, h)
    tbl = gf.table
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Emu(int(w * cw / total))
    tbl.rows[0].height = Inches(0.36)
    for ri in range(1, rows):
        tbl.rows[ri].height = Inches(row_h)
    for ri in range(rows):
        for ci in range(cols):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(data[ri][ci])
            if ri == 0:
                set_font(r, header_size, True, WHITE)
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
            else:
                set_font(r, fsize, False, GRAY)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 else RGBColor(0xED, 0xF1, 0xF7)
    return tbl

# =====================================================================
# P1 封面
# =====================================================================
s = new_slide()
add_rect(s, 0, 0, SW, SH, fill=NAVY, line=None, shape=MSO_SHAPE.RECTANGLE)
add_rect(s, 0, Inches(6.9), SW, Inches(0.6), fill=BLUE, line=None, shape=MSO_SHAPE.RECTANGLE)
add_text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.5),
         "数据治理 × 人工智能", size=16, bold=True, color=RGBColor(0x9D, 0xC3, 0xE6))
add_text(s, Inches(0.9), Inches(2.1), Inches(11.8), Inches(1.9),
         "利用 AI 进行数据治理\n全球案例与中国实践", size=40, bold=True, color=WHITE, line_spacing=1.15)
add_text(s, Inches(0.9), Inches(4.3), Inches(11.5), Inches(0.9),
         "从 Covestro 到魏桥创业集团：主数据治理的范式转移\nAI Agent / AI 原生 MDM / 多智能体协作平台", size=17, color=RGBColor(0xC9, 0xD9, 0xEE))
add_text(s, Inches(0.9), Inches(5.9), Inches(8), Inches(0.4),
         "汇报人：Frank（赵军军）    2026 年 8 月 25 日", size=14, color=RGBColor(0x9D, 0xC3, 0xE6))
add_text(s, Inches(9.0), Inches(5.9), Inches(3.5), Inches(0.4),
         "售前参考 · 技术分享", size=13, color=RGBColor(0x9D, 0xC3, 0xE6), align=PP_ALIGN.RIGHT)
add_notes(s, "开场：今天分享 AI 在数据治理领域的落地案例，重点是两个国内最新实践——魏桥创业集团的 AI 原生主数据管理，和用友 BIP 数据治理 Agents 协作平台，并对照国际标杆案例提炼共性方法。")

# =====================================================================
# P2 目录
# =====================================================================
s = new_slide()
add_slide_header(s, "CONTENTS", "目录", 2)
toc = [
    ("01", "为什么是现在", "AI 数据治理的驱动力与市场转折点"),
    ("02", "国际案例", "Covestro / Mastercard / Netflix 等五条主线全景"),
    ("03", "中国实践 ①", "魏桥创业集团 × 中翰软件：AI 原生主数据管理"),
    ("04", "中国实践 ②", "用友 BIP：数据治理 Agents 协作平台"),
    ("05", "启示与落地", "共性成功要素 + 项目落地路径建议"),
]
y = Inches(1.55)
for num, t, d in toc:
    add_rect(s, Inches(0.8), y, Inches(11.7), Inches(0.95), fill=BG)
    add_text(s, Inches(1.05), y + Inches(0.2), Inches(1.0), Inches(0.55), num, size=24, bold=True, color=BLUE)
    add_text(s, Inches(2.1), y + Inches(0.13), Inches(4.2), Inches(0.45), t, size=17, bold=True, color=NAVY)
    add_text(s, Inches(6.4), y + Inches(0.2), Inches(5.9), Inches(0.55), d, size=12.5, color=GRAY)
    y += Inches(1.1)
add_notes(s, "五部分：驱动力 → 国际案例 → 两个中国实践（魏桥、用友）→ 落地启示。")

# =====================================================================
# P3 驱动力：为什么是现在
# =====================================================================
s = new_slide()
add_slide_header(s, "01 · 为什么是现在", "五股力量把数据治理推向 AI 原生", 3)
cards = [
    ("数据孤岛与脏数据", "系统林立、标准不一，“烟囱式”数据各自为政；错误、重复、缺失的脏数据侵蚀决策质量——人工治理已到成本极限。", BLUE),
    ("“项目式”治理失效", "传统治理项目验收即停滞，“治理一次、沉睡多年”，规则更新滞后于业务变化，治理与业务“两张皮”。", ORANGE),
    ("政策转向“价值治理”", "从合规治理转向价值治理、从项目式转向常态化运营，DCMM/DAMA 标准体系化推进，治理成为企业标配能力。", GREEN),
    ("Agent 技术成熟", "2026 年大模型 Agent 从演示走向生产：多智能体协作、垂类大模型（如 BS-LM）、MCP 协议打通平台互操作。", BLUE),
    ("头部企业已跑通", "Covestro 主数据创建 12h→6min、Mastercard 元数据补全省 6,000+ 人时——AI 治理 ROI 已被验证。", GREEN),
]
x0, y0, cw, ch, gap = Inches(0.55), Inches(1.5), Inches(3.85), Inches(2.35), Inches(0.25)
for i, (t, d, ac) in enumerate(cards):
    cx = x0 + (cw + gap) * (i % 3)
    cy = y0 + (ch + Inches(0.28)) * (i // 3)
    add_card(s, cx, cy, cw, ch, t, d, accent=ac)
add_text(s, Inches(0.55), Inches(6.75), Inches(12.2), Inches(0.5),
         "结论：AI 不是给数据治理“加个功能”，而是把治理从“人工驱动”重构为“智能驱动”的拐点。",
         size=14, bold=True, color=NAVY)
add_notes(s, "引用用友 BIP 发布时的数据：自动化程度超 85%、效率提升 3 倍以上；百分点 BS-LM 垂类大模型基于近千政企项目语料。")

# =====================================================================
# P4 框架总览：五条主线
# =====================================================================
s = new_slide()
add_slide_header(s, "02 · 案例全景框架", "AI 参与数据治理的五条主线", 4)
main_lines = [
    ("① 主数据智能化", "AI Agent 直接参与物料/供应商/客户主数据的创建、查重、校验与变更审批", "Covestro MARIS/PARIS\n魏桥 AI 原生 MDM"),
    ("② 元数据自动补全", "基于血缘与目录的 AI 自动生成资产描述、归属、影响分析，替代人工编目", "Mastercard Context Agents\nNetflix DataHub"),
    ("③ 质量与合规监控", "敏感数据识别、分类分级、脱敏、质量规则自动执行与实时告警", "Banco do Brasil\n微软 Purview DSPM"),
    ("④ AI 用例自身治理", "对数百个 AI/GenAI 用例的登记、风险分级、审批与供应商证据管理", "Mastercard × Credo AI\nLSEG"),
    ("⑤ 数据资产产品化", "治理后的数据以“产品/市场”形态交付，内置契约与上下文，AI 可直接消费", "Mastercard Medallion\n麦当劳 × Collibra"),
]
y = Inches(1.5)
for name, desc, cases in main_lines:
    add_rect(s, Inches(0.55), y, Inches(12.2), Inches(0.92), fill=BG)
    add_text(s, Inches(0.8), y + Inches(0.1), Inches(2.7), Inches(0.7), name, size=15, bold=True, color=NAVY)
    add_text(s, Inches(3.6), y + Inches(0.12), Inches(5.6), Inches(0.7), desc, size=11.5, color=GRAY)
    add_text(s, Inches(9.4), y + Inches(0.12), Inches(3.2), Inches(0.7), cases, size=11, bold=True, color=BLUE)
    y += Inches(1.05)
add_notes(s, "五条主线对应上一轮调研的 8 个已验证国际案例；后文国际案例全景与两个中国案例都归入此框架。")

# =====================================================================
# P5 国际案例全景
# =====================================================================
s = new_slide()
add_slide_header(s, "02 · 国际案例", "已验证的国际案例全景（来源：厂商官方案例库）", 5)
data = [
    ["企业", "主线", "做法", "量化效果"],
    ["Covestro", "① 主数据", "AI Agent（MARIS/PARIS）+ SAP MDG，对话式创建、实时查重、外部地址核验、自动提交变更请求", "物料创建 12h → 6min（-99%），澄清请求 -99%"],
    ["Mastercard", "② 元数据", "Atlan 目录 + Context Agents 自动补全 3 万+ 资产元数据", "节省 6,000+ 人时"],
    ["Mastercard", "④ AI 用例", "Credo AI 注册表：数百 GenAI 用例风险分级、自动路由审批、供应商证据收集", "审批周期大幅缩短"],
    ["Netflix", "② 元数据", "DataHub 统一目录（数据/ML/软件资产），PII 覆盖率监控 + TTL 自动清理", "治理自服务化"],
    ["Banco do Brasil", "③ 合规", "IBM watsonx.governance + EY：AI 全生命周期监控（偏见/漂移/透明性）", "统一 AI 治理模型"],
    ["麦当劳", "⑤ 产品化", "Collibra 六支柱框架，95 个市场统一数据运营模型", "60 天建成治理底座"],
    ["LSEG", "④ AI 用例", "ChatGPT Enterprise + 治理嵌入（人审、模型评估、隐私控制）", "发布周期 3-6 月 → 2 周"],
]
add_table(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(4.6), data,
          col_w=[1.4, 1.1, 5.2, 3.1], fsize=10.5, row_h=0.42)
add_text(s, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.7),
         "注：另有全球制造业客户用微软 Purview DSPM for AI 在 AI 上线前完成敏感数据治理（AI-ready 前置模式）。",
         size=11.5, color=LGRAY)
add_notes(s, "全部数字来自 IBM/AWS/Atlan/DataHub/Credo AI/Collibra/OpenAI 官方案例页。Covestro 案例与本司 SAP MDG 主数据治理业务最相关。")

# =====================================================================
# P6 深度案例① Covestro
# =====================================================================
s = new_slide()
add_slide_header(s, "02 · 国际案例 ①", "Covestro：主数据创建从 12 小时压缩到 6 分钟", 6)
add_rect(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(1.3), fill=NAVY)
add_text(s, Inches(0.85), Inches(1.62), Inches(11.6), Inches(0.4),
         "背景：全球聚合物制造商（48 个生产基地 / 约 1.75 万员工）· AWS 官方博客披露 · Deloitte 实施", size=12.5, color=RGBColor(0xC9, 0xD9, 0xEE))
add_text(s, Inches(0.85), Inches(2.05), Inches(11.6), Inches(0.6),
         "业务用户创建物料主数据：原来靠邮件找 MDG 专家，平均 12 小时；现在与 AI Agent 对话，6 分钟完成", size=15.5, bold=True, color=WHITE)
kpis = [("12h → 6min", "主数据创建周期\n（-99%）"), ("~1,000 单/月", "物料请求量\n（约 1.2 万/年）"), ("99%", "澄清请求降幅\n（邮件往来消失）"), ("2 个 Agent", "MARIS（物料）\nPARIS（业务伙伴）")]
x = Inches(0.55)
for v, l in kpis:
    add_rect(s, x, Inches(3.05), Inches(2.9), Inches(1.25), fill=BG)
    add_text(s, x + Inches(0.15), Inches(3.18), Inches(2.6), Inches(0.5), v, size=22, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), Inches(3.72), Inches(2.6), Inches(0.5), l, size=10.5, color=GRAY, align=PP_ALIGN.CENTER)
    x += Inches(3.1)
add_text(s, Inches(0.55), Inches(4.55), Inches(12.2), Inches(0.4),
         "关键设计：AI 不直写生产库，全部通过 SAP MDG Change Request 走既有审批流", size=13, bold=True, color=NAVY)
add_card(s, Inches(0.55), Inches(5.05), Inches(6.0), Inches(1.9),
         "AI Agent 做了什么", "· 对话引导输入，只需少量关键字段\n· 跨系统实时查重（CDQ 数据镜像）\n· 地址与工商登记核验（D&B / BvD）\n· 自动补全属性（税号 VAT ID、地理编码）\n· 生成并提交 MDG 变更申请", accent=BLUE)
add_card(s, Inches(6.75), Inches(5.05), Inches(6.0), Inches(1.9),
         "业务价值", "· 自服务化：业务用户无需 SAP 技能\n· MDG 专家聚焦高价值治理任务\n· 全程可追溯，审计合规提升\n· 上线节奏：MARIS@ECC 2026 初；PARIS@S/4HANA 2027 初", accent=GREEN)
add_notes(s, "与我们的主数据治理方案对照：重复检测、地址校验、黄金记录补全三个 AI 注入点与五大域框架一一对应。架构：AWS Bedrock + SAP BTP/Integration Suite + SAP MDG。")

# =====================================================================
# P7 国际案例② Mastercard
# =====================================================================
s = new_slide()
add_slide_header(s, "02 · 国际案例 ②", "Mastercard：治理底座先行的“AI 就绪”样板", 7)
add_card(s, Inches(0.55), Inches(1.5), Inches(6.0), Inches(2.2),
         "十年治理积累 → AI 时代收割", "· BI 时代已建好：数据血缘、资产归属、元数据目录（Atlan）、Medallion 分层、数据产品市场\n· AI 到来时不必从零起步——治理即 AI 的“上下文底座”", accent=BLUE)
add_card(s, Inches(6.75), Inches(1.5), Inches(6.0), Inches(2.2),
         "Context Agents：元数据补全自动化", "· 基于既有血缘与目录，AI 自动为 3 万+ 资产生成描述、下游影响、消费方清单\n· 节省 6,000+ 人时的元数据引导工作\n· 通过 MCP 协议接入内部 Agent 工作流", accent=GREEN)
add_card(s, Inches(0.55), Inches(3.9), Inches(6.0), Inches(2.2),
         "AI 用例治理（× Credo AI）", "· 数百个 GenAI 用例统一登记、风险分级、自动路由审批\n· 供应商证据收集自动化（Vendor Portal）\n· 高管对 AI 使用全景可见", accent=ORANGE)
add_card(s, Inches(6.75), Inches(3.9), Inches(6.0), Inches(2.2),
         "对企业的启示", "· 分散式治理（联邦化）+ 中央底线控制\n· 数据契约与上下文内嵌于数据产品，而非事后补\n· 治理投入是复利：每新增 Agent 都复用同一底座", accent=NAVY)
add_notes(s, "Mastercard 案例说明：先有扎实治理底座，AI 才敢规模化。6,000 人时节省是多年目录/血缘投资的复利回报。")

# =====================================================================
# P8 中国视角：六平台格局
# =====================================================================
s = new_slide()
add_slide_header(s, "03 · 中国视角", "2026 年中国数据治理市场：进入 AI/Agent 时代", 8)
data = [
    ["厂商平台", "定位与差异化", "智能化深度", "开放/生态"],
    ["百分点 AI-DG", "数据治理垂类大模型 BS-LM（业内首个），对话式多智能体", "高：数万标准/规则/行业模型，集成效率 +80%、交付周期 -70%", "MCP 协议开放对接，支持私有化+第三方模型"],
    ["腾讯云 WeData", "Data+AI 一体化，DataOps+MLOps，首家通过信通院 DIOps", "中高：语义层 Unity Semantics + AI 助手（SQL 生成/纠错）", "MCP 目前偏语义层，生态偏腾讯系"],
    ["火山引擎 DataLeap", "分布式 + 多级治理，80+ 算子", "中：异常自动修复、元数据自动化", "2026 正式开源，云原生友好"],
    ["用友 BIP", "ERP 生态内“源头治理”，数据治理 Agents 协作平台（16 智能体）", "高：自动化率 85%+，治理效率 3 倍", "强绑定 BIP 生态，跨外部系统偏弱"],
    ["金蝶云·苍穹", "嵌入式 + 低代码建模，ERP 生态内治理", "中：低代码驱动", "绑定金蝶生态"],
    ["微软 Purview", "统一安全治理 + AI 全栈（DSPM for AI、Insider Risk）", "中高：AI 就绪前置", "Azure 生态强绑定"],
]
add_table(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(4.7), data,
          col_w=[1.6, 3.4, 3.4, 2.6], fsize=10, row_h=0.44)
add_text(s, Inches(0.55), Inches(6.5), Inches(12.2), Inches(0.6),
         "趋势：MCP 成为平台互操作标准——企业可保留现有数据基础设施，通过标准接口引入 AI 治理能力（IDC）。",
         size=12, bold=True, color=NAVY)
add_notes(s, "依据中国经济新闻网 2026-06-30《聚焦数据治理智能化：六家主流平台开放能力对比及企业选型指南》。")

# =====================================================================
# P9 中国实践① 魏桥：背景与痛点
# =====================================================================
s = new_slide()
add_slide_header(s, "04 · 中国实践 ①", "魏桥创业集团：AI 原生主数据管理（2026.08 签约）", 9)
add_rect(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(1.15), fill=NAVY)
add_text(s, Inches(0.85), Inches(1.62), Inches(11.6), Inches(0.9),
         "世界 500 强（2025 年第 166 位）· 山东第一大民企 · 连续 14 年入围 · 纺织 + 铝业两大主业",
         size=13.5, bold=True, color=WHITE)
facts = [("5,915 亿元", "2025 销售收入"), ("3,200 亿元", "总资产"), ("11 万人", "员工规模"), ("18 个", "国内外生产基地")]
x = Inches(0.55)
for v, l in facts:
    add_rect(s, x, Inches(2.85), Inches(2.9), Inches(1.1), fill=BG)
    add_text(s, x + Inches(0.15), Inches(2.95), Inches(2.6), Inches(0.5), v, size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), Inches(3.42), Inches(2.6), Inches(0.4), l, size=11, color=GRAY, align=PP_ALIGN.CENTER)
    x += Inches(3.1)
add_card(s, Inches(0.55), Inches(4.2), Inches(6.0), Inches(2.3),
         "数据治理痛点（多基地/多板块共性）", "· 主数据标准不统一：物料编码、客户/供应商信息对齐难\n· 跨系统协同效率不足：业务系统林立，接口靠人\n· 数据质量持久管控乏力：治理后容易“反弹”\n· 组织架构同步复杂：随业务单元增加呈指数级上升", accent=RED)
add_card(s, Inches(6.75), Inches(4.2), Inches(6.0), Inches(2.3),
         "战略选择：为何是“AI 原生”", "· 传统 MDM = 工具思维：表格录入、表单维护，高度依赖 IT 与外部顾问\n· 魏桥选择国产纯 AI 原生平台（中翰软件），私有化部署\n· 目标：业务人员自主管理数据，摆脱“项目结束、治理停滞”", accent=GREEN)
add_notes(s, "信息源：中国发展网/百家号多篇报道（2026-08-13/14）。魏桥是国内首个把纯 AI 原生 MDM 放入万亿级制造核心场景的企业。")

# =====================================================================
# P10 魏桥方案架构：Navigate OS
# =====================================================================
s = new_slide()
add_slide_header(s, "04 · 中国实践 ①", "魏桥方案：Navigate OS 底座上的纯 AI 原生 MDM", 10)
add_rect(s, Inches(0.55), Inches(1.45), Inches(12.2), Inches(0.75), fill=BG)
add_text(s, Inches(0.8), Inches(1.55), Inches(11.7), Inches(0.55),
         "中翰软件 · 国内首款纯 AI 原生主数据管理平台 · “一个对话框搞定一切主数据管理”",
         size=13, bold=True, color=NAVY)
layers = [
    ("交互层", "对话框即界面：拍照上传 / 复制粘贴 / 语音表达意图 → 智能查重 → 自动执行", BLUE),
    ("管理层", "主动式管理：推送待办、数据质量报告、清洗方案、分发失败原因、标准体系调整建议", GREEN),
    ("智能层", "三合一顾问（咨询+实施+技术）：自动生成项目计划/SOW、数据分类、编码结构、模型体系、接口脚本", ORANGE),
    ("底座层", "Navigate OS：智能工程（智能体可视化编排）+ 知识工程（知识库/领域模型）+ 本体论语义词典（数据推理）", NAVY),
]
y = Inches(2.35)
for name, desc, ac in layers:
    add_rect(s, Inches(0.55), y, Inches(12.2), Inches(0.85), fill=BG)
    add_rect(s, Inches(0.55), y, Inches(1.55), Inches(0.85), fill=ac, line=None, shape=MSO_SHAPE.RECTANGLE)
    add_text(s, Inches(0.55), y + Inches(0.24), Inches(1.55), Inches(0.4), name, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2.3), y + Inches(0.12), Inches(10.3), Inches(0.62), desc, size=12, color=GRAY)
    y += Inches(1.0)
add_text(s, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.6),
         "管理哲学转变：从“人适应系统”到“系统理解人”，从流程驱动到智能驱动，从“被动查询”到“主动预判”。",
         size=13, bold=True, color=NAVY)
add_notes(s, "Navigate OS 以管理者“顶视”为出发点；智能工程支持企业自主编排智能体，知识工程沉淀企业知识库。")

# =====================================================================
# P11 魏桥：传统 MDM vs AI 原生 MDM
# =====================================================================
s = new_slide()
add_slide_header(s, "04 · 中国实践 ①", "范式对比：传统 MDM vs AI 原生 MDM", 11)
data = [
    ["维度", "传统 MDM", "AI 原生 MDM（魏桥方案）"],
    ["交互方式", "菜单/表单/按钮，业务需培训", "对话框对话，拍照/粘贴/语音即可"],
    ["系统行为", "被动等待查询与录入", "主动推送待办、质量报告、清洗建议"],
    ["人员依赖", "高度依赖 IT 与外部顾问驻场", "业务人员自主管理，主导权回归甲方"],
    ["交付模式", "交钥匙工程，验收即撤离", "“双轨带练”：线上线下带教，能力内化"],
    ["数据能力", "规则匹配、查重靠人工/脚本", "本体论语义词典，数据具备推理与预判能力"],
    ["可持续性", "项目结束 → 治理停滞、质量反弹", "治理成为业务日常能力，长期成本下降"],
]
add_table(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(4.9), data,
          col_w=[1.7, 4.6, 5.9], fsize=11, row_h=0.42)
add_notes(s, "注意：以上差异基于厂商公开宣传（中翰软件官方发布），实际效果有待项目交付验证——魏桥文章亦有此保留（见“更大可能不等于必然实现”）。")

# =====================================================================
# P12 魏桥：行业意义
# =====================================================================
s = new_slide()
add_slide_header(s, "04 · 中国实践 ①", "魏桥案例的行业意义与借鉴点", 12)
cards = [
    ("国产纯 AI 原生首次进入万亿级制造", "同等体量制造企业此前几乎不用纯 AI 原生国产平台做核心主数据。魏桥选型本身具信号意义，可能改变国内数据治理市场产品方向。", BLUE),
    ("治理主导权交还业务", "核心机制：业务人员具备数据思维、成为真正数据管理者；治理从“一次性项目”变为“日常运营”。", GREEN),
    ("“带练”替代“交钥匙”", "线上线下协同带教，把数据管理能力沉淀到业务端，从源头维持数据质量，破解“治而不愈”顽疾。", ORANGE),
    ("风险与观察点（客观视角）", "万亿级业务场景对高并发、一致性、异常处理要求极高；AI 原生平台在生产环境的稳定性、人员适应速度、模型持续迭代均待验证。", RED),
]
x0, y0 = Inches(0.55), Inches(1.5)
for i, (t, d, ac) in enumerate(cards):
    cx = x0 + (i % 2) * Inches(6.25)
    cy = y0 + (i // 2) * Inches(2.6)
    add_card(s, cx, cy, Inches(6.0), Inches(2.4), t, d, accent=ac)
add_notes(s, "客观视角：文章（IDGA 署名）也提示“更大可能不等于必然实现”，具体效果要看产品成熟度与落地执行。")

# =====================================================================
# P13 中国实践② 用友 BIP 平台
# =====================================================================
s = new_slide()
add_slide_header(s, "04 · 中国实践 ②", "用友 BIP 数据治理 Agents 协作平台（2026.03 发布）", 13)
add_rect(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(1.1), fill=NAVY)
add_text(s, Inches(0.85), Inches(1.62), Inches(11.6), Inches(0.9),
         "用友 BIP 数智平台创新成果 · 首次将多智能体协作深度融入数据治理\n2026-03-16《用友 BIP 发布时刻》首期直播正式发布",
         size=13, bold=True, color=WHITE)
kpis = [("16 个", "首批专业智能体\n五大类 · 超百项技能"), ("85%+", "治理流程自动化程度"), ("90%+", "人工重复劳动替代率"), ("3 倍+", "治理效率提升")]
x = Inches(0.55)
for v, l in kpis:
    add_rect(s, x, Inches(2.8), Inches(2.9), Inches(1.25), fill=BG)
    add_text(s, x + Inches(0.15), Inches(2.92), Inches(2.6), Inches(0.5), v, size=21, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), Inches(3.45), Inches(2.6), Inches(0.55), l, size=10.5, color=GRAY, align=PP_ALIGN.CENTER)
    x += Inches(3.1)
add_text(s, Inches(0.55), Inches(4.35), Inches(12.2), Inches(0.4),
         "五大智能体类别：自然语言处理 · 业务架构建模 · 数据建模管理 · 数据价值转化 · 数据治理核心能力",
         size=13, bold=True, color=NAVY)
add_card(s, Inches(0.55), Inches(4.85), Inches(6.0), Inches(2.0),
         "“人机共治”范式", "· 覆盖业务调研→架构梳理→标准设计→质量规则→湖仓落标→常态化运营全流程\n· 在线设计、在线协同、在线运营一体化\n· 智能体自动化 + 人工精准校验", accent=BLUE)
add_card(s, Inches(6.75), Inches(4.85), Inches(6.0), Inches(2.0),
         "技术可信设计（白盒）", "· 所有智能体操作过程、决策逻辑可追溯可查看，解决大模型幻觉\n· 技能基于 DAMA、DCMM 等标准封装，适配行业数据字典\n· 研发前走访多家行业标杆，先内部试用再推广", accent=GREEN)
add_notes(s, "来源：CSDN 转载用友官方稿 + 雪球。自动化程度 85%+ 为厂商发布口径。")

# =====================================================================
# P14 用友：价值与市场定位
# =====================================================================
s = new_slide()
add_slide_header(s, "04 · 中国实践 ②", "用友 BIP：从“一次性工程”到“可持续运营”", 14)
add_card(s, Inches(0.55), Inches(1.5), Inches(6.0), Inches(2.4),
         "解决的核心矛盾", "· 传统治理“靠人堆、靠加班、靠救火”，项目结束治理即停\n· 治理规则更新滞后于业务变化，形成“治理一次、沉睡多年”\n· 数据治理与业务发展“两张皮”", accent=RED)
add_card(s, Inches(6.75), Inches(1.5), Inches(6.0), Inches(2.4),
         "成本模型重构", "· 从“线性投入”（持续堆人力）\n→ “一次性智能体投入 + 长期少量线性投入”\n· 降低治理门槛，中小资源企业也能系统性治理\n· 让数据治理从“专业少数人专属”变为企业“标配”", accent=GREEN)
add_card(s, Inches(0.55), Inches(4.1), Inches(6.0), Inches(2.3),
         "市场定位（六平台对比）", "· “ERP 生态内源头治理”——治理能力与 BIP 业务系统强绑定\n· 优势：业务系统内数据直接治理，形成闭环\n· 局限：跨外部系统/异构生态的适配范围偏窄", accent=BLUE)
add_card(s, Inches(6.75), Inches(4.1), Inches(6.0), Inches(2.3),
         "对企业决策的含义", "· 已有用友 BIP 生态的企业：治理可低摩擦嵌入现有系统\n· 异构系统为主的企业：需评估 MCP/开放对接能力或组合多家平台\n· 选型本质 = 生态绑定深度 × 智能化能力 × 开放性 的权衡", accent=NAVY)
add_notes(s, "对比来源：经济新闻网六平台测评——用友 BIP 平台 Agent 协作+源头治理，强绑定 BIP 生态，跨外部系统偏弱。")

# =====================================================================
# P15 六平台横向对比总结
# =====================================================================
s = new_slide()
add_slide_header(s, "04 · 横向对比", "智能化深度 × 平台开放性：选型坐标", 15)
add_rect(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(4.6), fill=WHITE, line=BORDER)
# 坐标轴
add_text(s, Inches(0.9), Inches(1.62), Inches(11.5), Inches(0.35), "智能化深度（AI/Agent 能力）→", size=12, bold=True, color=LGRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.62), Inches(3.0), Inches(0.35), Inches(2.2), "开放性\n↑", size=11, bold=True, color=LGRAY)
positions = [
    ("百分点 AI-DG", Inches(3.4), Inches(2.0), BLUE),
    ("用友 BIP",       Inches(5.6), Inches(2.3), ORANGE),
    ("腾讯 WeData",    Inches(4.6), Inches(3.6), BLUE),
    ("火山 DataLeap",  Inches(6.6), Inches(4.4), BLUE),
    ("金蝶苍穹",       Inches(8.6), Inches(3.0), GRAY),
    ("微软 Purview",   Inches(9.8), Inches(2.2), GRAY),
]
for name, bx, by, col in positions:
    add_rect(s, bx, by, Inches(2.6), Inches(0.5), fill=col, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, bx, by + Inches(0.1), Inches(2.6), Inches(0.3), name, size=11.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.8),
         "选型建议：生态匹配优先（在用友/金蝶生态内选 BIP/苍穹，云原生选 WeData/DataLeap，AI 深度优先选 AI-DG）\n再横向比较：MCP 开放对接能力将是跨平台组合的关键（IDC：Agent 标准协议降低集成难度）",
         size=12, bold=True, color=NAVY)
add_notes(s, "坐标为定性判断（基于六平台测评），供选型讨论用；正式选型需按客户系统现状打分。")

# =====================================================================
# P16 共性成功要素
# =====================================================================
s = new_slide()
add_slide_header(s, "05 · 启示", "跨案例的六条共性成功要素", 16)
elements = [
    ("AI 推荐 + 人审批", "无一例外：AI 不直写生产库。Covestro 走 MDG 审批流，用友“智能体自动化+人工校验”，魏桥对话确认机制", BLUE),
    ("主导权交还业务", "魏桥“带练”交付、Covestro 自服务化：业务人员成为数据管理者，治理才可持续", GREEN),
    ("治理底座先行", "Mastercard 十年目录/血缘/契约积累 → AI 复利；微软 Purview 客户“AI 就绪前置”", NAVY),
    ("技能标准化封装", "用友按 DAMA/DCMM 封装智能体技能；中翰领域模型+语义词典——治理知识产品化", ORANGE),
    ("白盒可解释", "用友白盒设计解决幻觉信任；IBM watsonx 全程可追溯——AI 治理本身要被治理", BLUE),
    ("从项目走向运营", "成本模型从线性人力投入转向“一次性平台+少量线性”，治理常态化是共同终点", GREEN),
]
y = Inches(1.5)
for t, d, ac in elements:
    add_rect(s, Inches(0.55), y, Inches(12.2), Inches(0.78), fill=BG)
    add_rect(s, Inches(0.55), y, Inches(2.6), Inches(0.78), fill=ac, line=None, shape=MSO_SHAPE.RECTANGLE)
    add_text(s, Inches(0.55), y + Inches(0.19), Inches(2.6), Inches(0.4), t, size=13.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(3.35), y + Inches(0.1), Inches(9.25), Inches(0.6), d, size=11.5, color=GRAY)
    y += Inches(0.9)
add_notes(s, "六条要素贯穿所有案例。讲售前时重点用 Covestro 审批流 + 魏桥带练模式，回答客户“AI 直接改数据安全吗”的顾虑。")

# =====================================================================
# P17 落地路径建议
# =====================================================================
s = new_slide()
add_slide_header(s, "05 · 启示", "对我司/客户项目的落地路径建议（PoC 三步走）", 17)
steps = [
    ("第一步 · 选域与基线（2-4 周）", "选定 1 个主数据域（物料或供应商）做试点\n建立治理收益基线：重复数据量+关联采购金额、必填字段完整率、跨系统映射覆盖率、工单闭环时长", BLUE),
    ("第二步 · AI 注入点 PoC（4-8 周）", "对照 Covestro 三个注入点：AI 查重/外部数据核验/黄金记录补全\n方案 = AI 推荐 + 人工审批 + 可控执行 + 可回滚\n输出：查重准确率、创建周期缩短、审批通过率", GREEN),
    ("第三步 · 运营化与扩展（3-6 月）", "“带练”式交接：业务人员掌握日常数据管理\n智能体技能按 DCMM/DAMA 封装，沉淀为可复用资产\n扩展至 BOM/工艺路线/客户域，接入 MCP 与外部系统", ORANGE),
]
x = Inches(0.55)
for t, d, ac in steps:
    add_rect(s, x, Inches(1.5), Inches(3.95), Inches(4.3), fill=BG)
    add_rect(s, x, Inches(1.5), Inches(3.95), Inches(0.55), fill=ac, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, x + Inches(0.15), Inches(1.58), Inches(3.65), Inches(0.45), t, size=12.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(2.25), Inches(3.55), Inches(3.4), d, size=11, color=GRAY)
    x += Inches(4.15)
add_text(s, Inches(0.55), Inches(6.15), Inches(12.2), Inches(0.9),
         "PoC 验收铁律：每个验收标准必须绑定业务影响指标（如：查重准确率 → 重复采购金额下降）——治理项目收益不可见是最大落地障碍。",
         size=13, bold=True, color=NAVY)
add_notes(s, "基线五指标来自本司数据治理方法论（治理收益基线）；三步走与 PoC 验收标准可直接用于售前方案。")

# =====================================================================
# P18 风险与对策
# =====================================================================
s = new_slide()
add_slide_header(s, "05 · 启示", "风险与对策：把话说在前面的售前姿态", 18)
data = [
    ["风险", "表现", "对策"],
    ["AI 效果被高估", "厂商宣传口径（85% 自动化、3 倍效率）多为发布稿数据，缺少独立第三方验证", "PoC 用客户真实数据复测；指标以实测为准（参考魏桥报道“更大可能≠必然实现”）"],
    ["生态绑定陷阱", "用友/金蝶平台与自有生态强绑定，跨异构系统能力偏弱", "先做系统现状盘点；用 MCP/开放 API 评估组合方案，不把赌注押单一平台"],
    ["数据安全与合规", "AI 访问敏感数据（PII/税号/财务数据）存在泄露与越权风险", "沿用“AI 不直写、全程留痕、人审把关”；参考 Purview DSPM、BB 全生命周期监控"],
    ["组织抗拒", "数据治理是组织权力再分配，业务部门可能消极配合", "Kickoff 前置条件：发起人 VP+ 书面授权；“带练”式交付降低业务抵触"],
    ["可持续运营缺失", "项目结束治理停滞的顽疾在新方案下依然存在", "把“运营化”写进合同交付物：双轨带练、技能资产沉淀、月度健康度指标"],
]
add_table(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(4.8), data,
          col_w=[1.8, 5.0, 5.2], fsize=10.5, row_h=0.45)
add_notes(s, "售前姿态：不回避风险，用 PoC 实测+合同化运营承诺建立信任。")

# =====================================================================
# P19 信息来源与置信度
# =====================================================================
s = new_slide()
add_slide_header(s, "附 · 信息源", "案例数据来源与置信度说明", 19)
data = [
    ["案例", "关键数据", "来源", "置信度"],
    ["Covestro", "12h→6min；~1,000 单/月；澄清 -99%", "AWS for SAP 官方博客（AWS 官方披露）", "高 [KNOWN]"],
    ["Mastercard", "30,000+ 资产；6,000+ 人时", "Atlan 客户案例页（厂商披露）", "高 [KNOWN]"],
    ["Netflix/麦当劳/BB/LSEG", "定性成果", "DataHub/Collibra/IBM/OpenAI 案例页", "中高 [KNOWN]"],
    ["魏桥创业集团", "500强166位；5,915亿；3,200亿；11万人；18基地；2026.08 签约", "中国发展网/百家号多篇（中国发展改革署名）", "高 [KNOWN]"],
    ["中翰 AI 原生 MDM", "Navigate OS；首款纯 AI 原生；双轨带练", "同花顺财经/中国发展网（厂商发布稿）", "中 [KNOWN]（厂商口径）"],
    ["用友 BIP Agents", "2026.03.16 发布；16 智能体；85%+；3 倍", "用友官方发布稿（CSDN/雪球转载）", "中 [KNOWN]（厂商口径）"],
    ["六平台对比", "智能化深度/开放性定性评价", "中国经济新闻网 2026-06-30 测评", "中 [KNOWN]"],
]
add_table(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(4.8), data,
          col_w=[2.0, 3.2, 3.6, 3.0], fsize=10, row_h=0.42)
add_text(s, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.6),
         "口径说明：国内案例量化数字均来自厂商/媒体发布稿，未经验证为“宣传口径”；售前引用时建议标注并 PoC 实测。",
         size=11.5, color=LGRAY)
add_notes(s, "诚实标注口径，售前与分享时保持一致，避免被客户追问数字来源时失分。")

# =====================================================================
# P20 总结
# =====================================================================
s = new_slide()
add_slide_header(s, "SUMMARY", "三个关键信息", 20)
keys = [
    ("①", "范式已切换", "AI 数据治理从“演示”走向“生产”：Covestro 12h→6min、用友 16 智能体 85% 自动化、魏桥把纯 AI 原生 MDM 放进万亿级制造核心场景"),
    ("②", "路径有共识", "AI 推荐+人审批 · 治理底座先行 · 主导权交还业务 · 白盒可解释 · 从项目走向运营——六条共性跨中外案例一致"),
    ("③", "落地靠实测", "厂商宣传口径 ≠ 客户实测效果：PoC 绑定业务指标、带练式交付、运营化写入合同，才能让治理真正持续"),
]
y = Inches(1.6)
for n, t, d in keys:
    add_rect(s, Inches(0.55), y, Inches(12.2), Inches(1.35), fill=BG)
    add_text(s, Inches(0.85), y + Inches(0.3), Inches(0.8), Inches(0.7), n, size=30, bold=True, color=BLUE)
    add_text(s, Inches(1.75), y + Inches(0.18), Inches(2.4), Inches(0.5), t, size=18, bold=True, color=NAVY)
    add_text(s, Inches(4.3), y + Inches(0.2), Inches(8.2), Inches(1.0), d, size=12, color=GRAY)
    y += Inches(1.55)
add_rect(s, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.6), fill=NAVY)
add_text(s, Inches(0.55), Inches(6.67), Inches(12.2), Inches(0.4),
         "Q & A  —  谢谢", size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_notes(s, "收尾：三个关键信息 + 提问环节。可附 PoC 三步走（P17）作为后续讨论钩子。")

# ---------- 保存 ----------
out = r"D:\AI\14 - 数据治理\AI数据治理案例与趋势_20260825.pptx"
prs.save(out)
print("saved:", out, "| slides:", len(prs.slides.__iter__.__self__._sldIdLst))
