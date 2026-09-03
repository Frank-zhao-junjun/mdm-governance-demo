# -*- coding: utf-8 -*-
"""
构建《AI 在主数据治理中的应用(2).pptx》
- 16:9, 13.333 x 7.5 英寸, 微软雅黑, 深蓝 #1F4E79 + 青色 #2E9E97/#00B0A8
- 事实口径纪律:
  * 中翰软件(Navigate OS/对话式交互/双轨带练/国内首款) -> 标注"据中翰软件官方发布"
  * 用友效果数字(85%/90%/3倍) -> 标注"用友官方称"
  * 魏桥 5915亿/18基地/11万员工 可直引
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- 常量 ----------
DARK = RGBColor(0x1F, 0x4E, 0x79)      # 深蓝主色
TEAL = RGBColor(0x2E, 0x9E, 0x97)      # 青色点缀
TEAL2 = RGBColor(0x00, 0xB0, 0xA8)     # 亮青
GRAY_BG = RGBColor(0xF2, 0xF4, 0xF7)   # 浅灰卡片底
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x33, 0x33, 0x33)      # 正文深灰
MUTE = RGBColor(0x8A, 0x94, 0xA6)      # 弱化灰
LINE = RGBColor(0xD9, 0xDE, 0xE4)      # 分隔线
CARD_LINE = RGBColor(0xE0, 0xE4, 0xEA)
MID_BLUE = RGBColor(0x3A, 0x6E, 0xA5)  # 中间蓝
SLATE = RGBColor(0x5B, 0x6B, 0x7F)     # 灰蓝
FONT = "微软雅黑"

PAGE_W, PAGE_H = 13.333, 7.5
MARGIN = 0.5
CONTENT_W = PAGE_W - 2 * MARGIN  # 12.333
FOOTER = "AI Native 数据治理方案 · 2026-08"

prs = Presentation()
prs.slide_width = Inches(PAGE_W)
prs.slide_height = Inches(PAGE_H)
BLANK = prs.slide_layouts[6]


# ---------- 基础工具 ----------
def _set_run(run, size, bold, color, font=FONT):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    f.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", font)


def add_box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def add_para(tf, text, size=12, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
             first=False, space_after=4, space_before=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    _set_run(run, size, bold, color)
    return p


def add_rect(slide, x, y, w, h, fill, line_color=None, shape=MSO_SHAPE.RECTANGLE, radius=None):
    shp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    if radius is not None:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    return shp


def shape_text(shp, lines, anchor=MSO_ANCHOR.MIDDLE, m=0.12):
    """lines: list of (text, size, bold, color, align)"""
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(m)
    tf.margin_right = Inches(m)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    for i, (text, size, bold, color, align) in enumerate(lines):
        add_para(tf, text, size, bold, color, align, first=(i == 0), space_after=2)


def add_card(slide, x, y, w, h, title, body_lines, fill=GRAY_BG,
             title_color=DARK, title_size=14, body_size=11.5, accent=True):
    add_rect(slide, x, y, w, h, fill, CARD_LINE, MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.055)
    if accent:
        add_rect(slide, x + 0.06, y + 0.16, 0.055, h - 0.32, TEAL)
    tb, tf = add_box(slide, x + 0.22, y + 0.12, w - 0.36, h - 0.24)
    add_para(tf, title, title_size, True, title_color, first=True, space_after=6)
    for line in body_lines:
        add_para(tf, line, body_size, False, TEXT, space_after=3)


def add_header(slide, title, page_no):
    add_rect(slide, MARGIN, 0.40, 0.09, 0.44, TEAL)
    tb, tf = add_box(slide, 0.74, 0.30, 11.6, 0.64)
    add_para(tf, title, 24, True, DARK, first=True)
    add_rect(slide, MARGIN, 1.02, CONTENT_W, 0.014, LINE)
    # 页脚
    _, ff = add_box(slide, MARGIN, 7.12, 8.0, 0.3)
    add_para(ff, FOOTER, 9, False, MUTE, first=True, space_after=0)
    _, pf = add_box(slide, 11.8, 7.12, 1.0, 0.3)
    add_para(pf, str(page_no), 9, False, MUTE, PP_ALIGN.RIGHT, first=True, space_after=0)


def add_banner(slide, text, y=6.0, h=0.62, fill=DARK, size=13.5):
    shp = add_rect(slide, MARGIN, y, CONTENT_W, h, fill, None,
                   MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.18)
    shape_text(shp, [(text, size, True, WHITE, PP_ALIGN.CENTER)])
    return shp


def add_note(slide, text, y, color=MUTE, size=9, align=PP_ALIGN.RIGHT):
    _, tf = add_box(slide, MARGIN, y, CONTENT_W, 0.3)
    add_para(tf, text, size, False, color, align, first=True, space_after=0)


def add_table(slide, x, y, w, data, col_widths, row_h=0.52, header_h=0.5,
              header_size=12.5, body_size=11.5):
    rows, cols = len(data), len(data[0])
    total_h = header_h + (rows - 1) * row_h
    gf = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(total_h))
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = Inches(cw)
    for r in range(rows):
        tbl.rows[r].height = Inches(header_h if r == 0 else row_h)
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = DARK
            else:
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else GRAY_BG
            tf = cell.text_frame
            tf.word_wrap = True
            add_para(tf, str(data[r][c]),
                     header_size if r == 0 else body_size,
                     r == 0,
                     WHITE if r == 0 else TEXT,
                     PP_ALIGN.LEFT, first=True, space_after=0)
    return gf


def new_slide(title=None, page_no=None, notes=None):
    s = prs.slides.add_slide(BLANK)
    if title is not None:
        add_header(s, title, page_no)
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


# ---------- 第1页 封面 ----------
s = new_slide(notes="开场：本方案回答一个问题——AI 如何重塑制造业主数据治理。先看两个标杆，再给出我们的方案。")
add_rect(s, 0, 0, PAGE_W, PAGE_H, DARK)
add_rect(s, 0, 0, PAGE_W, 0.09, TEAL2)
add_rect(s, 1.6, 2.42, 0.9, 0.07, TEAL2)
_, tf = add_box(s, 1.6, 2.62, 10.1, 1.1)
add_para(tf, "AI 在主数据治理中的应用（二）", 32, True, WHITE, first=True)
_, tf = add_box(s, 1.6, 3.72, 10.1, 0.6)
add_para(tf, "AI Native 数据治理方案（制造业）", 18, False, RGBColor(0xBF, 0xD7, 0xEA), first=True)
_, tf = add_box(s, 1.6, 6.35, 10.5, 0.4)
add_para(tf, "借鉴中翰软件 Navigate OS × 用友 BIP 数据治理 Agents 协作平台 · 2026-08",
         11, False, RGBColor(0x9F, 0xB3, 0xC8), first=True)

# ---------- 第2页 为什么需要 AI Native 数据治理 ----------
s = new_slide("为什么需要 AI Native 数据治理", 2,
              notes="四个痛点相互强化：越依赖专家和人工，越难持续。行业共识正在形成——人机共治是方向。")
_, tf = add_box(s, MARGIN, 1.18, CONTENT_W, 0.35)
add_para(tf, "传统主数据治理模式的四个结构性痛点：", 13, True, TEXT, first=True)
pains = [
    ("高成本", ["依赖专职团队与长周期项目，", "投入大、回报慢"]),
    ("高壁垒", ["依赖稀缺的数据治理专家，", "知识难以沉淀与复制"]),
    ("低效率", ["人工录入、核对、分发，", "耗时且易出错"]),
    ("难持续", ["项目式治理，上线即退化，", "成果难以保持"]),
]
cw, gap = 2.95, 0.177
for i, (t, lines) in enumerate(pains):
    add_card(s, MARGIN + i * (cw + gap), 1.62, cw, 3.6, t, lines,
             title_size=16, body_size=12)
add_banner(s, "行业拐点：从「人工治理」走向「人机共治」——AI 正从辅助工具变为治理的操作系统",
           y=5.62, h=0.66)

# ---------- 第3页 标杆借鉴① 中翰 ----------
s = new_slide("标杆借鉴 ①：中翰软件 AI 原生 MDM（Navigate OS）", 3,
              notes="中翰走 AI 原生、自下而上路线。平台能力描述均来自厂商官方发布，尚无第三方实测。")
add_card(s, MARGIN, 1.16, CONTENT_W, 0.98, "案例背景",
         ["魏桥创业集团（2025 年营收 5915 亿元、18 个生产基地、11 万员工），2026 年 8 月与中翰软件签署主数据治理合作。"],
         title_size=13, body_size=11.5)
feats = [
    ("对话式交互", ["一个对话框完成全部操作；", "拍照上传或复制粘贴即可新增 / 变更，", "平台智能查重并自动执行"]),
    ("主动智能", ["主动推送数据质量报告与清洗方案；", "主动分析分发失败原因，", "建议治理体系调整"]),
    ("多模型融合底座", ["融合多模态大模型、领域模型、", "小模型、本体模型与知识库，", "各司其职"]),
    ("双轨带练实施法", ["线上线下协同赋能，", "带练业务人员上手操作，", "让治理能力沉淀在业务侧"]),
]
cw, ch, gx, gy = 6.05, 1.86, 0.233, 0.16
for i, (t, lines) in enumerate(feats):
    x = MARGIN + (i % 2) * (cw + gx)
    y = 2.32 + (i // 2) * (ch + gy)
    add_card(s, x, y, cw, ch, t, lines, title_size=14, body_size=11)
add_note(s, "注：本页平台特征描述均据中翰软件官方发布，暂无第三方实测佐证。", 6.42)

# ---------- 第4页 标杆借鉴② 用友 ----------
s = new_slide("标杆借鉴 ②：用友 BIP 数据治理 Agents 协作平台", 4,
              notes="用友走平台化、自上而下路线。85%/90%/3 倍为厂商自报指标，引用时注意口径。")
_, tf = add_box(s, MARGIN, 1.14, CONTENT_W, 0.4)
add_para(tf, "2026 年 3 月 16 日发布 · 定位：数据治理从「项目式」转向「常态化运营」 · 首批 16 个专业智能体，分五大类",
         12.5, True, TEXT, first=True)
tbl_data = [
    ("智能体类别", "核心能力"),
    ("自然语言处理类", "理解业务语言，将业务描述转化为治理规则"),
    ("业务架构建模类", "梳理业务架构，建立治理框架"),
    ("数据建模管理类", "设计数据模型，管理数据资产"),
    ("数据价值转化类", "释放数据价值，支撑业务决策"),
    ("数据治理核心能力类", "执行治理全流程，保障数据质量"),
]
add_table(s, MARGIN, 1.66, 7.3, tbl_data, [2.3, 5.0], row_h=0.62, header_h=0.52)
rx = 8.06
rw = PAGE_W - MARGIN - rx  # 4.773
add_card(s, rx, 1.66, rw, 2.42, "关键指标（用友官方称）",
         ["· 自动化程度超 85%",
          "· 替代 90% 以上人工重复劳动",
          "· 治理效率提升 3 倍以上",
          "（以上为厂商自报指标，暂无第三方验证）"],
         title_size=14, body_size=12)
add_card(s, rx, 4.28, rw, 1.62, "白盒设计",
         ["所有智能体的操作过程、决策逻辑", "均可追溯、可查看，拒绝黑箱操作。"],
         title_size=14, body_size=12)
add_note(s, "注：发布时间与智能体分类据用友官网；效果数字须冠以「用友官方称」。", 6.42)

# ---------- 第5页 两条路径对比与启示 ----------
s = new_slide("两条路径对比与启示", 5,
              notes="两条路线不是二选一：底座学用友的体系化，切入学中翰的场景化。")
add_card(s, MARGIN, 1.35, 6.05, 4.15, "自下而上 · 中翰路线（场景驱动）",
         ["· 业务场景驱动，小切口切入",
          "· AI 原生对话式交互",
          "· 业务人员自主操作，无需专家",
          "· 快速见效，易于推广复制"],
         title_size=15, body_size=13)
add_card(s, 6.783, 1.35, 6.05, 4.15, "自上而下 · 用友路线（平台驱动）",
         ["· 平台化、体系化能力建设",
          "· 多智能体协同作业",
          "· 标准化能力底座，可复用",
          "· 常态化运营，持续演进"],
         title_size=15, body_size=13)
add_banner(s, "启示：制造业最优解 = 平台底座（自上而下）× 场景切入（自下而上），双轮驱动",
           y=5.85, h=0.66)

# ---------- 第6页 方案总体架构（五层） ----------
s = new_slide("方案总体架构：五层架构", 6,
              notes="读法：L5 集成现状系统，L4 是主数据服务本体，L3/L2 是 AI 增量，L1 面向业务用户。")
layers = [
    ("L1 交互层", "对话式治理门户 · 移动端拍照录入 · IM 集成", TEAL2),
    ("L2 智能体层", "主数据治理智能体集群（标准 / 建模 / 查重 / 质量 / 分发 / 影响分析 / 问答）", TEAL),
    ("L3 AI 能力层", "多模态大模型 · 领域模型 · 小模型 · 本体模型 · 知识库 · 向量库", DARK),
    ("L4 主数据服务层", "数据标准 · 数据建模 · 主数据管理 · 质量监控 · 分发订阅", MID_BLUE),
    ("L5 集成层", "ERP（SAP / 金蝶 / 用友）· MES · WMS · PLM · SRM", SLATE),
]
ly, lh, lg = 1.24, 0.98, 0.135
for i, (label, content, color) in enumerate(layers):
    y = ly + i * (lh + lg)
    shp = add_rect(s, MARGIN, y, CONTENT_W, lh, color, None,
                   MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run(); r1.text = label + "　"
    _set_run(r1, 14.5, True, WHITE)
    r2 = p.add_run(); r2.text = content
    _set_run(r2, 12, False, WHITE)
add_note(s, "注：L2/L3 为 AI 增量能力，L4/L5 与现有主数据平台及业务系统对接。", 6.85, align=PP_ALIGN.LEFT)

# ---------- 第7页 核心设计① 对话式治理入口 ----------
s = new_slide("核心设计 ①：对话式治理入口", 7,
              notes="对话式入口是降低使用门槛的关键，理念借鉴中翰（厂商官方口径），按制造业场景落地。")
_, tf = add_box(s, MARGIN, 1.2, CONTENT_W, 0.4)
add_para(tf, "一个对话框，完成主数据新增 / 变更 / 查询 / 审批全流程", 14, True, DARK, first=True)
caps = [
    ("多模态录入", ["拍照识别图纸、铭牌信息；", "粘贴文本智能解析，", "自动结构化为标准字段"]),
    ("智能查重", ["新增即查重，", "实时提示相似历史物料，", "从源头防止重复建档"]),
    ("填单辅助", ["自动补全物料属性，", "推荐分类与编码，", "大幅降低填写门槛"]),
]
cw, gap = 3.95, 0.24
for i, (t, lines) in enumerate(caps):
    add_card(s, MARGIN + i * (cw + gap), 1.85, cw, 3.4, t, lines,
             title_size=15, body_size=12.5)
add_banner(s, "设计理念借鉴中翰软件对话式交互（据中翰软件官方发布），结合制造业场景重构",
           y=5.6, h=0.62, fill=TEAL, size=12.5)

# ---------- 第8页 核心设计② 智能体集群 ----------
s = new_slide("核心设计 ②：制造业治理智能体集群", 8,
              notes="7 个智能体覆盖治理全流程，强调实时协同——不是串行接力，而是并行响应、相互校验。")
_, tf = add_box(s, MARGIN, 1.14, CONTENT_W, 0.4)
add_para(tf, "参考用友五大类智能体划分，映射为制造业 7 个专业智能体：", 13, True, TEXT, first=True)
agents = [
    ("智能体", "职责"),
    ("标准管理 Agent", "数据标准生成与维护"),
    ("数据建模 Agent", "属性模板设计与管理"),
    ("智能查重 Agent", "重复 / 近似物料识别"),
    ("质量监控 Agent", "数据异常检测与预警"),
    ("分发编排 Agent", "新增 / 变更 / 冻结事件驱动分发"),
    ("影响分析 Agent", "变更影响分析（BOM / 工艺 / 库存 / 订单）"),
    ("治理问答 Agent", "治理规则问答与流程指引"),
]
add_table(s, MARGIN, 1.62, CONTENT_W, agents, [3.2, 9.133], row_h=0.5, header_h=0.48,
          header_size=12.5, body_size=11.5)
add_banner(s, "关键特征：多智能体实时协同，而非串行接力——一次请求，多 Agent 并行响应、相互校验",
           y=6.1, h=0.6, fill=TEAL, size=12.5)

# ---------- 第9页 核心设计③ 多模型融合底座 ----------
s = new_slide("核心设计 ③：多模型融合技术底座", 9,
              notes="多模型各司其职是成本与效果平衡的关键：不是所有任务都值得用大模型。")
comps_top = [
    ("多模态大模型", ["理解图纸、铭牌照片与", "自然语言 —— 负责「理解与交互」"]),
    ("领域模型", ["沉淀物料分类体系与命名规范", "—— 负责「行业 know-how」"]),
    ("小模型", ["相似度查重、异常检测，", "低延迟低成本 —— 负责「高频判断」"]),
]
comps_bottom = [
    ("本体模型", ["构建物料 - BOM - 工艺 - 供应商", "关系图谱 —— 负责「关系推理」"]),
    ("知识库 + 向量库", ["标准文档、历史案例检索增强（RAG）", "—— 负责「知识供给」"]),
]
cw, gap = 3.95, 0.24
for i, (t, lines) in enumerate(comps_top):
    add_card(s, MARGIN + i * (cw + gap), 1.3, cw, 2.1, t, lines, title_size=14.5, body_size=11.5)
cw2, gap2 = 6.05, 0.233
for i, (t, lines) in enumerate(comps_bottom):
    add_card(s, MARGIN + i * (cw2 + gap2), 3.6, cw2, 2.1, t, lines, title_size=14.5, body_size=11.5)
add_banner(s, "分工原则：大模型做理解与交互，小模型做高频判断，本体做关系推理——各司其职",
           y=6.0, h=0.62)

# ---------- 第10页 核心设计④ 主动式数据质量管理 ----------
s = new_slide("核心设计 ④：主动式数据质量管理", 10,
              notes="主动式质量管理的要点：AI 发起、人工确认、规则回流，形成自我强化的闭环。")
steps = [
    ("1 监控", "持续监控\n质量指标"),
    ("2 预警", "主动推送\n质量报告"),
    ("3 诊断", "定位分发\n失败原因"),
    ("4 建议", "生成清洗\n方案"),
    ("5 执行", "人工确认后\n自动清洗"),
    ("6 复核", "效果评估\n规则回流"),
]
bw, bh, agap = 1.83, 1.75, 0.2
by = 1.55
for i, (t, d) in enumerate(steps):
    x = MARGIN + i * (bw + agap)
    blk = add_rect(s, x, by, bw, bh, DARK if i % 2 == 0 else TEAL, None,
                   MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.09)
    tf = blk.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_para(tf, t, 14, True, WHITE, PP_ALIGN.CENTER, first=True, space_after=6)
    for seg in d.split("\n"):
        add_para(tf, seg, 10.5, False, WHITE, PP_ALIGN.CENTER, space_after=1)
    if i < 5:
        ar = add_rect(s, x + bw + 0.015, by + bh / 2 - 0.11, 0.17, 0.22, TEAL2, None,
                      MSO_SHAPE.RIGHT_ARROW)
add_card(s, MARGIN, 3.75, CONTENT_W, 1.05, "闭环要点",
         ["预警与诊断由 AI 主动发起；清洗执行前必须人工确认；复核结果回流，持续优化监控规则。"],
         title_size=13, body_size=11.5)
add_banner(s, "从「被动查询」到「主动预判」——数据质量管理成为持续运转的闭环",
           y=5.35, h=0.62)

# ---------- 第11页 落地场景 ----------
s = new_slide("制造业落地场景：六大治理对象 + 三个小切口", 11,
              notes="治理对象与前期主数据治理专题一致；试点选小切口：见效快、可量化、易推广。")
_, tf = add_box(s, MARGIN, 1.14, CONTENT_W, 0.35)
add_para(tf, "六大治理对象（与主数据治理专题范围一致）", 13, True, DARK, first=True)
objs = ["物料", "供应商", "客户", "BOM", "工艺路线", "员工"]
ow, ogap = 1.93, 0.15
for i, name in enumerate(objs):
    shp = add_rect(s, MARGIN + i * (ow + ogap), 1.56, ow, 0.95, GRAY_BG, CARD_LINE,
                   MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    shape_text(shp, [(name, 14.5, True, TEAL, PP_ALIGN.CENTER)])
_, tf = add_box(s, MARGIN, 2.85, CONTENT_W, 0.35)
add_para(tf, "试点建议：三个小切口（见效快 · 可量化 · 易推广）", 13, True, DARK, first=True)
pilots = [
    ("重复物料检测", ["对存量物料库智能查重，输出", "重复 / 近似清单，", "直接支撑采购整合与降本"]),
    ("分类与属性推荐", ["新增物料自动推荐分类、", "属性与编码，提升一致性，", "降低业务人员填写门槛"]),
    ("数据质量预警", ["异常检测与主动预警推送，", "问题早发现、早处置，", "防止脏数据向下游扩散"]),
]
cw, gap = 3.95, 0.24
for i, (t, lines) in enumerate(pilots):
    x = MARGIN + i * (cw + gap)
    add_card(s, x, 3.3, cw, 2.5, t, lines, title_size=14.5, body_size=11.5)
    badge = add_rect(s, x + cw - 1.15, 3.3 + 0.14, 0.95, 0.3, TEAL2, None,
                     MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    shape_text(badge, [("优先试点", 9.5, True, WHITE, PP_ALIGN.CENTER)], m=0.02)

# ---------- 第12页 可信与可控 ----------
s = new_slide("可信与可控：白盒 + 人机协同", 12,
              notes="可信可控是 AI 治理能落地的底线：白盒、人在回路、权限审计，缺一不可。")
guards = [
    ("白盒可追溯", ["智能体操作过程与决策逻辑全留痕，",
                  "可追溯、可查看；借鉴用友白盒设计思想，",
                  "拒绝黑箱操作。"]),
    ("人在回路", ["AI 建议、人决策：",
                "新增、冻结、批量清洗等关键节点，",
                "必须人工复核确认后方可生效。"]),
    ("权限与审计", ["分级权限管理、操作审计日志、",
                  "合规留痕，",
                  "满足内控与外部审计要求。"]),
]
cw, gap = 3.95, 0.24
for i, (t, lines) in enumerate(guards):
    add_card(s, MARGIN + i * (cw + gap), 1.5, cw, 3.5, t, lines,
             title_size=15, body_size=12.5)
add_banner(s, "原则：AI 提效，人类把关；过程透明，责任清晰", y=5.5, h=0.62)

# ---------- 第13页 实施路线图 ----------
s = new_slide("实施路线图：三阶段推进", 13,
              notes="三阶段节奏：先小切口验证价值，再扩展对象与系统，最后转入常态化运营。")
phases = [
    ("试点期（1–3 个月）", TEAL2,
     ["· 小切口场景上线", "  （查重 / 推荐 / 预警）",
      "· 建立数据标准与质量基线",
      "· 对话式入口试运行"]),
    ("推广期（3–9 个月）", TEAL,
     ["· 扩展至六大治理对象",
      "· 接入 ERP / MES / WMS", "  / PLM / SRM",
      "· 智能体集群全面启用"]),
    ("运营期（9 个月 +）", DARK,
     ["· 治理常态化运营",
      "· 「双轨带练」赋能业务人员", "  自主治理",
      "· 持续优化模型与规则"]),
]
cw, gap = 4.0, 0.166
for i, (t, color, lines) in enumerate(phases):
    x = MARGIN + i * (cw + gap)
    ch_shape = add_rect(s, x, 1.45, cw, 0.62, color, None, MSO_SHAPE.CHEVRON)
    try:
        ch_shape.adjustments[0] = 0.35
    except Exception:
        pass
    shape_text(ch_shape, [(t, 14, True, WHITE, PP_ALIGN.CENTER)])
    add_card(s, x, 2.35, cw - 0.1, 3.35, t.split("（")[0], lines,
             title_size=13.5, body_size=11.5)
add_note(s, "注：「双轨带练」为中翰软件提出的实施方法论（据中翰软件官方发布）。", 6.35)

# ---------- 第14页 价值度量与结语 ----------
s = new_slide("价值度量与结语", 14,
              notes="价值用五维框架表达，量化指标先留白、试点期后回填实测值。以结语收束。")
vals = [
    ("效率", "建档与变更周期缩短"),
    ("成本", "重复采购与呆滞库存下降"),
    ("质量", "准确率 / 完整率 / 一致率提升"),
    ("风控", "变更可追溯、可审计"),
    ("战略", "夯实 AI 应用的数据底座"),
]
vw, vgap = 2.33, 0.17
for i, (t, d) in enumerate(vals):
    x = MARGIN + i * (vw + vgap)
    add_card(s, x, 1.25, vw, 1.85, t, [d], title_size=15, body_size=11)
add_card(s, MARGIN, 3.35, CONTENT_W, 1.15, "量化指标模板（试点期后回填实测值）",
         ["重复物料下降 xx%　｜　建档周期缩短 xx%　｜　变更同步及时率 xx%　｜　数据质量得分提升 xx 分"],
         title_size=13, body_size=12.5)
add_rect(s, 5.42, 5.05, 2.5, 0.05, TEAL2)
_, tf = add_box(s, MARGIN, 5.3, CONTENT_W, 0.9)
add_para(tf, "让数据治理的主导权交还业务，让数据从「负担」变为「资产」。",
         20, True, DARK, PP_ALIGN.CENTER, first=True)

# ---------- 保存与验证 ----------
OUT = r"e:/14 - 数据治理/AI 在主数据治理中的应用(2).pptx"
prs.save(OUT)
print("saved:", OUT)

prs2 = Presentation(OUT)
W, H = prs2.slide_width, prs2.slide_height
print("slides:", len(prs2.slides))
TOL = Emu(9525)  # 0.01 英寸容差
issues = []
for i, sl in enumerate(prs2.slides, 1):
    print(f"slide {i}: {len(sl.shapes)} shapes")
    for sh in sl.shapes:
        l, t, w, h = sh.left, sh.top, sh.width, sh.height
        if l is None:
            continue
        if l < -TOL or t < -TOL or l + w > W + TOL or t + h > H + TOL:
            issues.append((i, sh.shape_id, sh.name, int(l), int(t), int(w), int(h)))
print("bounds issues:", issues if issues else "NONE")
