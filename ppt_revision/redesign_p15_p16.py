# -*- coding: utf-8 -*-
"""
重新设计 P15 / P16：
- 删除老旧的五边形流程、3D圆柱、连接线
- 换成现代编号圆圈+连线步骤条
- 插入 HTML 渲染的高保真 UI 截图
- 统一配色：#1B3A5C 深蓝 / #2E75B6 蓝 / #2BA89C 青 / #475569 灰
"""
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy, os

SRC = r"D:\AI\14 - 数据治理\ppt_revision\AI_主数据治理_v8.pptx"
DST = r"D:\AI\14 - 数据治理\ppt_revision\AI_主数据治理_v9.pptx"
IMG = r"D:\AI\14 - 数据治理\ppt_revision\new_images"

# 配色
C_DARK   = RGBColor(0x1B, 0x3A, 0x5C)
C_BLUE   = RGBColor(0x2E, 0x75, 0xB6)
C_TEAL   = RGBColor(0x2B, 0xA8, 0x9C)
C_GRAY   = RGBColor(0x47, 0x55, 0x69)
C_LGRAY  = RGBColor(0x94, 0xA3, 0xB8)
C_BG     = RGBColor(0xF1, 0xF5, 0xF9)
C_LINE   = RGBColor(0xCB, 0xD5, 0xE1)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation(SRC)

def set_font(run, size=11, color=C_GRAY, bold=False, name="微软雅黑"):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', name)

def add_textbox(slide, x, y, w, h, text, size=11, color=C_GRAY, bold=False,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Cm(0.05)
    tf.margin_top = tf.margin_bottom = Cm(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_font(r, size, color, bold)
    return tb

def add_circle(slide, cx, cy, d, fill=C_BLUE, line_color=None):
    """以 (cx,cy) 为圆心画圆，直径 d(cm)"""
    x = cx - d/2
    y = cy - d/2
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(x), Cm(y), Cm(d), Cm(d))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    return sh

def add_rounded_rect(slide, x, y, w, h, fill=C_BG, line_color=C_LINE, text="",
                     text_size=10, text_color=C_GRAY, bold=False):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    if text:
        tf = sh.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Cm(0.1)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        set_font(r, text_size, text_color, bold)
    # 调整圆角
    sh.adjustments[0] = 0.3
    return sh

def add_line(slide, x1, y1, x2, y2, color=C_LINE, width=1.5):
    from pptx.enum.shapes import MSO_CONNECTOR
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    return conn

def add_arrow(slide, x1, y1, x2, y2, color=C_LGRAY, width=1.5):
    from pptx.enum.shapes import MSO_CONNECTOR
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    # 添加箭头
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    return conn

def clear_slide_keep(slide, keep_keywords):
    """删除所有形状，除了文本包含指定关键词的（标题和场景说明）"""
    to_remove = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            txt = sh.text_frame.text
            if any(kw in txt for kw in keep_keywords):
                continue
        to_remove.append(sh)
    for sh in to_remove:
        sh._element.getparent().remove(sh._element)

def restyle_title(slide, keyword, size=24):
    """统一标题样式"""
    for sh in slide.shapes:
        if sh.has_text_frame and keyword in sh.text_frame.text:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    set_font(r, size, C_DARK, True)
            return sh
    return None

# ============================================================
# P15: 场景3 — 智能查数&创建新数据
# ============================================================
slide15 = prs.slides[14]
clear_slide_keep(slide15, ["场景3", "场景说明"])
restyle_title(slide15, "场景3")

# 场景说明样式统一
for sh in slide15.shapes:
    if sh.has_text_frame and "场景说明" in sh.text_frame.text:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                set_font(r, 11, C_GRAY, False)
        sh.left = Cm(1.5); sh.top = Cm(1.6); sh.width = Cm(30); sh.height = Cm(1.2)

# --- 现代步骤条 ---
step_y = 3.6  # 圆心 y
steps = ["多模态识别", "智能匹配", "填写编码申请", "MDM规范性检查", "提交审核"]
x_start = 2.2
spacing = 5.6
circle_d = 0.85

# 连接线（先画线，圆圈盖在上面）
add_line(slide15, x_start + circle_d/2, step_y,
         x_start + spacing*(len(steps)-1) - circle_d/2, step_y, C_LINE, 1.5)

for i, label in enumerate(steps):
    cx = x_start + spacing * i
    color = C_TEAL if i % 2 == 0 else C_BLUE
    circ = add_circle(slide15, cx, step_y, circle_d, fill=color)
    # 编号
    tf = circ.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(i+1)
    set_font(r, 11, C_WHITE, True)
    # 标签
    add_textbox(slide15, cx - 2.0, step_y + circle_d/2 + 0.15, 4.0, 0.8,
                label, size=10, color=C_GRAY, align=PP_ALIGN.CENTER)

# 箭头 + OA/ERP 端点
end_x = x_start + spacing*(len(steps)-1) + 1.8
add_arrow(slide15, x_start + spacing*(len(steps)-1) + circle_d/2 + 0.2, step_y,
          end_x - 0.1, step_y, C_LGRAY, 1.5)
add_rounded_rect(slide15, end_x, step_y - 0.55, 3.0, 1.1, fill=C_BG, line_color=C_LINE,
                 text="OA / ERP", text_size=11, text_color=C_GRAY, bold=True)

# --- UI 截图 ---
chat_path = os.path.join(IMG, "p15_chat.png")
form_path = os.path.join(IMG, "p15_form.png")

# 左：聊天界面
chat_w = 15.0
chat_h = chat_w * 520 / 640
chat_x = 2.2
chat_y = 5.4
slide15.shapes.add_picture(chat_path, Cm(chat_x), Cm(chat_y), Cm(chat_w), Cm(chat_h))

# 右：表单
form_w = 11.5
form_h = form_w * 520 / 520  # 1:1
form_x = chat_x + chat_w + 1.5
form_y = 5.4
slide15.shapes.add_picture(form_path, Cm(form_x), Cm(form_y), Cm(form_w), Cm(form_h))

# 小标注
add_textbox(slide15, chat_x, chat_y + chat_h + 0.15, chat_w, 0.6,
            "拍照 / 粘贴 / 语音输入，AI 识别铭牌并匹配相似编码",
            size=9, color=C_LGRAY, align=PP_ALIGN.CENTER)
add_textbox(slide15, form_x, form_y + form_h + 0.15, form_w, 0.6,
            "AI 自动填充字段，人工核对后提交",
            size=9, color=C_LGRAY, align=PP_ALIGN.CENTER)

print("P15 redesigned")

# ============================================================
# P16: 场景4 — 主动管理
# ============================================================
slide16 = prs.slides[15]
clear_slide_keep(slide16, ["场景4", "场景说明"])
restyle_title(slide16, "场景4")

for sh in slide16.shapes:
    if sh.has_text_frame and "场景说明" in sh.text_frame.text:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                set_font(r, 11, C_GRAY, False)
        sh.left = Cm(1.5); sh.top = Cm(1.6); sh.width = Cm(30); sh.height = Cm(1.2)

# --- 左：告警中心截图 ---
alerts_path = os.path.join(IMG, "p16_alerts.png")
alerts_w = 19.5
alerts_h = alerts_w * 500 / 760
alerts_x = 1.5
alerts_y = 3.3
slide16.shapes.add_picture(alerts_path, Cm(alerts_x), Cm(alerts_y), Cm(alerts_w), Cm(alerts_h))

# --- 右：能力清单 ---
right_x = 22.5
right_w = 10.0
cap_y = 3.6
cap_d = 0.9
cap_spacing = 1.85

capabilities = [
    ("推送待办", "异常和待办主动推送给责任人"),
    ("质量报告", "数据质量日报、周报自动生成"),
    ("清洗方案", "AI 发现问题并给出清洗建议"),
    ("分发告警", "分发失败即时告警并附原因"),
    ("标准建议", "标准体系调整建议，待人工评审"),
]

# 竖向连接线
add_line(slide16, right_x + cap_d/2, cap_y + cap_d,
         right_x + cap_d/2, cap_y + cap_spacing*(len(capabilities)-1),
         C_LINE, 1.5)

for i, (title, desc) in enumerate(capabilities):
    cy = cap_y + cap_spacing * i
    color = C_TEAL if i % 2 == 0 else C_BLUE
    circ = add_circle(slide16, right_x + cap_d/2, cy + cap_d/2, cap_d, fill=color)
    tf = circ.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(i+1)
    set_font(r, 11, C_WHITE, True)
    # 标题
    add_textbox(slide16, right_x + cap_d + 0.3, cy, right_w - cap_d - 0.3, 0.7,
                title, size=12, color=C_DARK, bold=True)
    # 描述
    add_textbox(slide16, right_x + cap_d + 0.3, cy + 0.65, right_w - cap_d - 0.3, 0.8,
                desc, size=9, color=C_LGRAY)

# 箭头 + 推送渠道端点
end_y = cap_y + cap_spacing*(len(capabilities)-1) + 1.3
add_arrow(slide16, right_x + cap_d/2, cap_y + cap_spacing*(len(capabilities)-1) + cap_d + 0.1,
          right_x + cap_d/2, end_y - 0.15, C_LGRAY, 1.5)
add_rounded_rect(slide16, right_x - 0.3, end_y, right_w + 0.6, 1.0,
                 fill=C_BG, line_color=C_LINE,
                 text="邮件 / IM / 站内信推送", text_size=10, text_color=C_GRAY, bold=True)

print("P16 redesigned")

prs.save(DST)
print(f"saved -> {DST}")
