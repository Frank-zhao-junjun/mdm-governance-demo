# -*- coding: utf-8 -*-
from pptx import Presentation

path = r"D:\AI\14 - 数据治理\ppt_revision\AI_主数据治理_v3.pptx"
prs = Presentation(path)

targets = [1, 6, 11, 12, 14, 18, 19, 20, 21, 22, 23, 24, 25]
for idx in targets:
    slide = prs.slides[idx-1]
    print(f"\n{'='*60}\n### 第 {idx} 页\n{'='*60}")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text
        if txt.strip():
            print(f"\n[{shape.name}]")
            print(txt)
