# -*- coding: utf-8 -*-
"""
v4：在 v3 基础上，按“干过的实干家”标准再过一遍
- 去口号、大词、排比，换成具体动作、踩过的坑、取舍边界
- 占位指标改成“待实测”，不装数字
- 同步改部分备注
"""
from pptx import Presentation

SRC = r"D:\AI\14 - 数据治理\ppt_revision\AI_主数据治理_v3.pptx"
DST = r"D:\AI\14 - 数据治理\ppt_revision\AI_主数据治理_v4.pptx"

prs = Presentation(SRC)


def set_para_text(para, text):
    runs = list(para.runs)
    if runs:
        runs[0].text = text
        for r in runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        run = para.add_run()
        run.text = text


def find_shape(slide, name):
    for sh in slide.shapes:
        if sh.name == name and sh.has_text_frame:
            return sh
    raise ValueError(f"shape not found: {name}")


def edit(slide_num, shape_name, sets=None):
    slide = prs.slides[slide_num - 1]
    shape = find_shape(slide, shape_name)
    paras = list(shape.text_frame.paragraphs)
    if sets:
        for idx, txt in sets.items():
            set_para_text(paras[idx], txt)


def set_notes(slide_num, text):
    slide = prs.slides[slide_num - 1]
    slide.notes_slide.notes_text_frame.text = text


# ---------- 第 1 页：封面备注 ----------
set_notes(1, "开场：本方案回答一个问题——AI 在制造业主数据治理里到底能干什么。先看中某油和 SAP MDG 两个实践，再给我们的方案。")

# ---------- 第 12 页：AI 平台总览，实干口吻重写 ----------
edit(12, "TextBox 2", sets={
    0: "查重、校验、清洗这些重复活交给 AI，人只在拿不准和要拍板时介入",
})
edit(12, "TextBox 11", sets={
    1: "查询、新增、变更，都在一个对话框；",
    2: "铭牌拍照或粘贴文本，自动识别填单；",
    3: "提交前列出相似编码，避免重复建档",
})
edit(12, "TextBox 14", sets={
    1: "先学存量数据的规律，定出合理阈值；",
    2: "质量报告按时推，不用人催；",
    3: "发现异常先记下来，附清洗建议；",
    4: "可疑数据主动预警；",
    5: "分发失败自动查原因，不用翻日志",
})
edit(12, "TextBox 17", sets={
    1: "大模型看懂图纸、铭牌和自然语言；",
    2: "知识库放标准规范，本体模型管物料-BOM-工艺-供应商关系",
})
edit(12, "TextBox 20", sets={
    1: "清洗过的脏数据、驳回理由都攒下来，更新阈值和案例；",
    2: "标准变了，知识库跟着更新；",
    3: "常见问题不用重教，AI 按上次规则处理",
})
set_notes(12, "这页讲清 AI 接手什么、人保留什么：重复规则判断交给 AI，取舍和拍板留给人。")

# ---------- 第 18 页：对话式入口 ----------
edit(18, "TextBox 6", sets={
    0: "查编码、提申请、改属性在一个对话框完成，审批仍走现有流程",
})
edit(18, "TextBox 9", sets={
    3: "自动填到对应字段里",
})
edit(18, "TextBox 12", sets={
    2: "把相似历史编码列出来，",
})

# ---------- 第 19 页：智能体集群 ----------
edit(19, "TextBox 6", sets={
    0: "7 个智能体分工负责标准、建模、查重、质量、分发、影响分析和问答，由 AI 助手统一调度。",
})
edit(19, "Rounded Rectangle 8", sets={
    0: "关键特征：一次请求可触发多个 Agent 并行处理、互相校验，不用一个传一个地等",
})
set_notes(19, "7 个智能体覆盖治理全流程。协同方式是并行、互相校验，不是串行接力。")

# ---------- 第 20 页：多模型底座 ----------
edit(20, "TextBox 11", sets={
    1: "固化物料分类体系和命名规范，分类编码不靠个人经验",
})

# ---------- 第 21 页：主动式质量 ----------
edit(21, "TextBox 19", sets={
    1: "预警和诊断由 AI 主动发起；清洗方案必须人确认后才执行；复核结果用来调整监控阈值和规则，不用人手动改。",
})
set_notes(21, "要点三条：AI 发起预警和诊断，清洗必须人确认，复核结果用来调规则。")

# ---------- 第 22 页：落地场景 ----------
edit(22, "TextBox 16", sets={
    3: "结果直接给采购合并同款、减少重复采购",
})
edit(22, "TextBox 20", sets={
    3: "业务人员不用死记字段和编码规则",
})
set_notes(22, "治理对象与前期专题一致。试点从三个小切口进，先做出可量化的效果再推广。")

# ---------- 第 23 页：可信可控备注 ----------
set_notes(23, "可信可控是落地底线：白盒、人在回路、权限审计，哪一块缺了都不敢上线。")

# ---------- 第 24 页：价值度量 ----------
edit(24, "TextBox 8", sets={
    2: "建档周期缩短率（待实测）",
    3: "变更同步及时率（待实测）",
})
edit(24, "TextBox 11", sets={
    2: "重复物料下降率（待实测）",
})
edit(24, "TextBox 14", sets={
    2: "数据质量得分提升（待实测）",
})
edit(24, "TextBox 20", sets={
    1: "把主数据管准，后面的 AI 应用才站得住",
})
set_notes(24, "五维价值。量化指标待试点实测后回填，不拍脑袋。")

# ---------- 第 25 页：结尾 ----------
edit(25, "文本框 2", sets={
    0: "把主数据管准，业务才跑得顺",
})

prs.save(DST)
print(f"saved -> {DST}")
