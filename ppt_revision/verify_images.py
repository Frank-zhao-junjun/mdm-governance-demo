# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
import os

path = r"D:\AI\14 - 数据治理\ppt_revision\AI_主数据治理_v5.pptx"
prs = Presentation(path)
print("slides:", len(prs.slides))
out = r"D:\AI\14 - 数据治理\ppt_revision\verify_media"
os.makedirs(out, exist_ok=True)

for idx in [15, 16]:
    slide = prs.slides[idx-1]
    rels = slide.part.rels
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            blip = sh._element.findall('.//' + qn('a:blip'))[0]
            rid = blip.get(qn('r:embed'))
            part = rels[rid].target_part
            ext = part.content_type.split('/')[-1].replace('jpeg', 'jpg')
            safe_name = sh.name.replace(' ', '')
            fn = os.path.join(out, f"p{idx}_{safe_name}.{ext}")
            with open(fn, 'wb') as f:
                f.write(part.blob)
            print(f"P{idx} {sh.name}: {len(part.blob)} bytes -> {fn}")
