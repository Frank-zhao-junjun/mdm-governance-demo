# -*- coding: utf-8 -*-
from pptx import Presentation

path = r"D:\AI\14 - 数据治理\ppt_revision\AI_主数据治理_v4.pptx"
prs = Presentation(path)

# 抽查关键页 + 确认文件完整
print(f"总页数: {len(prs.slides)}")
for idx in [12, 18, 19, 20, 21, 22, 23, 24, 25]:
    slide = prs.slides[idx-1]
    print(f"\n{'='*60}\n### 第 {idx} 页\n{'='*60}")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text
        if txt.strip():
            print(f"\n[{shape.name}]")
            print(txt)
    if slide.has_notes_slide:
        n = slide.notes_slide.notes_text_frame.text
        if n.strip():
            print(f"\n[备注] {n}")
