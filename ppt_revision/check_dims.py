# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Emu

path = r"D:\AI\14 - 数据治理\ppt_revision\AI_主数据治理_v3.pptx"
prs = Presentation(path)

def emu_cm(v):
    return round(v / 360000, 2)

checks = [
    (12, "TextBox 2"),
    (12, "TextBox 11"),
    (12, "TextBox 14"),
    (12, "TextBox 17"),
    (12, "TextBox 20"),
    (18, "TextBox 6"),
    (19, "TextBox 6"),
    (19, "Rounded Rectangle 8"),
    (20, "TextBox 8"),
    (20, "TextBox 11"),
    (20, "TextBox 14"),
    (20, "TextBox 17"),
    (20, "TextBox 20"),
    (20, "Rounded Rectangle 21"),
    (21, "TextBox 19"),
    (21, "Rounded Rectangle 20"),
    (23, "TextBox 8"),
    (23, "TextBox 11"),
    (23, "TextBox 14"),
    (23, "Rounded Rectangle 15"),
    (24, "TextBox 8"),
    (24, "TextBox 11"),
    (24, "TextBox 14"),
    (24, "TextBox 17"),
    (24, "TextBox 20"),
    (25, "文本框 2"),
    (22, "TextBox 16"),
    (22, "TextBox 20"),
]
for sn, name in checks:
    slide = prs.slides[sn-1]
    for sh in slide.shapes:
        if sh.name == name:
            w = emu_cm(sh.width); h = emu_cm(sh.height)
            # word wrap
            try:
                wrap = sh.text_frame.word_wrap
            except Exception:
                wrap = "?"
            print(f"P{sn} {name:28s} {w}cm x {h}cm  wrap={wrap}")
