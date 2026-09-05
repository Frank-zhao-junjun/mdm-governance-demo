# -*- coding: utf-8 -*-
"""
替换 P15/P16 的图片：
- P15 图片3 (image21) -> p15_chat.png   对话式查数界面
- P15 图片27 (image25) -> p15_form.png   AI填好的物料表单
- P16 图片3 (image21) -> p16_alerts.png  主动告警中心
通过替换 blip 关系实现，保留形状位置/尺寸/图层。
"""
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

SRC = r"D:\AI\14 - 数据治理\ppt_revision\AI_主数据治理_v4.pptx"
DST = r"D:\AI\14 - 数据治理\ppt_revision\AI_主数据治理_v5.pptx"
IMG_DIR = r"D:\AI\14 - 数据治理\ppt_revision\new_images"

prs = Presentation(SRC)

def replace_picture_image(slide, shape_name, image_path):
    """找到指定图片形状，把它的 blip 指向新图片（新增 image part + 新关系）。"""
    for sh in slide.shapes:
        if sh.name != shape_name:
            continue
        blips = sh._element.findall('.//' + qn('a:blip'))
        if not blips:
            raise RuntimeError(f"{shape_name}: no blip found")
        # 新增图片 part（挂在 package 上）
        image_part = prs.part.package.get_or_add_image_part(image_path)
        # 建立 slide part -> image part 的关系
        new_rId = slide.part.relate_to(image_part, RT.IMAGE)
        for blip in blips:
            blip.set(qn('r:embed'), new_rId)
        print(f"  replaced {shape_name} -> {image_path} (rId={new_rId})")
        return
    raise RuntimeError(f"shape not found: {shape_name}")

print("第 15 页:")
replace_picture_image(prs.slides[14], "图片 3",  IMG_DIR + r"\p15_chat.png")
replace_picture_image(prs.slides[14], "图片 27", IMG_DIR + r"\p15_form.png")

print("第 16 页:")
replace_picture_image(prs.slides[15], "图片 3",  IMG_DIR + r"\p16_alerts.png")

prs.save(DST)
print(f"\nsaved -> {DST}")
