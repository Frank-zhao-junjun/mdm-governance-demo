# -*- coding: utf-8 -*-
"""v4 收尾：收紧第12页两个窄框，防止文字溢出"""
from pptx import Presentation

PATH = r"D:\AI\14 - 数据治理\ppt_revision\AI_主数据治理_v4.pptx"
prs = Presentation(PATH)

def set_para_text(para, text):
    runs = list(para.runs)
    if runs:
        runs[0].text = text
        for r in runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        para.add_run().text = text

def edit(sn, name, sets):
    slide = prs.slides[sn-1]
    for sh in slide.shapes:
        if sh.name == name and sh.has_text_frame:
            paras = list(sh.text_frame.paragraphs)
            for i, t in sets.items():
                set_para_text(paras[i], t)

# 多模型融合底座：两句各压到一行
edit(12, "TextBox 17", {
    1: "大模型负责理解，知识库提供标准规范；",
    2: "本体模型管物料-BOM-工艺-供应商关系",
})
# 知识沉淀：第一句压到一行
edit(12, "TextBox 20", {
    1: "脏数据、驳回理由都攒下来更新案例；",
})

prs.save(PATH)
print("v4 收紧完成")

# 复核这两块最终文本
slide = prs.slides[11]
for sh in slide.shapes:
    if sh.name in ("TextBox 17", "TextBox 20"):
        print(f"\n[{sh.name}]")
        print(sh.text_frame.text)
