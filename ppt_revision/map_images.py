# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

path = r"D:\AI\14 - 数据治理\ppt_revision\AI_主数据治理_v4.pptx"
prs = Presentation(path)

for idx in [15, 16]:
    slide = prs.slides[idx-1]
    print(f"\n### 第 {idx} 页")
    # map rId -> target
    rels = slide.part.rels
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # get blip rId
            blip = sh._element.findall('.//' + qn('a:blip'))
            for b in blip:
                rid = b.get(qn('r:embed'))
                if rid and rid in rels:
                    target = rels[rid].target_ref
                    print(f"  shape='{sh.name}' pos=({sh.left},{sh.top}) size=({sh.width},{sh.height}) -> {target}")
