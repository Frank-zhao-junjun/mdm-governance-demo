# -*- coding: utf-8 -*-
"""修复 v7 的两个硬伤，输出 v8"""
from pptx import Presentation

SRC = r"D:\AI\14 - 数据治理\ppt_revision\AI_主数据治理_v7.pptx"
DST = r"D:\AI\14 - 数据治理\ppt_revision\AI_主数据治理_v8.pptx"

prs = Presentation(SRC)

def set_para_text(para, text):
    runs = list(para.runs)
    if runs:
        runs[0].text = text
        for r in runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        para.add_run().text = text

def remove_para(para):
    para._p.getparent().remove(para._p)

# ---- P15: 删除文本框4的第三段残留垃圾 ----
slide15 = prs.slides[14]
for sh in slide15.shapes:
    if sh.name == "文本框 4":
        paras = list(sh.text_frame.paragraphs)
        print("P15 文本框4 段落:")
        for i, p in enumerate(paras):
            print(f"  P{i}: {repr(p.text)}")
        # 删除最后一段（残留的 |-----> 告诉物料编号、名称。）
        remove_para(paras[-1])
        print("  -> 已删除最后一段")
        break

# ---- P23: 修复白盒可追溯卡片的逗号孤行 ----
slide23 = prs.slides[22]
for sh in slide23.shapes:
    if sh.name == "TextBox 8":
        paras = list(sh.text_frame.paragraphs)
        print("\nP23 TextBox8 段落:")
        for i, p in enumerate(paras):
            print(f"  P{i}: {repr(p.text)}")
        set_para_text(paras[1], "智能体每一步操作和判断依据都留痕，可查看、可回溯；")
        set_para_text(paras[2], "借鉴用友白盒思路，不给黑箱结论。")
        remove_para(paras[3])
        print("  -> 已重排并删除多余段")
        break

prs.save(DST)
print(f"\nsaved -> {DST}")
