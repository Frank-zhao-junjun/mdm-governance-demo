# -*- coding: utf-8 -*-
"""下载重生成的两张图并替换 P15 的图片，输出 v6"""
import urllib.request, os
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

IMG_DIR = r"D:\AI\14 - 数据治理\ppt_revision\new_images"
SRC = r"D:\AI\14 - 数据治理\ppt_revision\AI_主数据治理_v5.pptx"
DST = r"D:\AI\14 - 数据治理\ppt_revision\AI_主数据治理_v6.pptx"

urls = {
    "p15_chat.png": "https://aka.doubaocdn.com/s/FTHhVaIG2k",
    "p15_form.png": "https://aka.doubaocdn.com/s/JQ2FpXTOiB",
}
for name, url in urls.items():
    path = os.path.join(IMG_DIR, name)
    req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
    with urllib.request.urlopen(req, timeout=60) as r, open(path, "wb") as f:
        f.write(r.read())
    print("downloaded", name, os.path.getsize(path))

prs = Presentation(SRC)

def replace_picture_image(slide, shape_name, image_path):
    for sh in slide.shapes:
        if sh.name != shape_name:
            continue
        blip = sh._element.findall('.//' + qn('a:blip'))[0]
        image_part = prs.part.package.get_or_add_image_part(image_path)
        new_rId = slide.part.relate_to(image_part, RT.IMAGE)
        blip.set(qn('r:embed'), new_rId)
        print(f"replaced {shape_name} -> {os.path.basename(image_path)}")
        return
    raise RuntimeError(shape_name)

replace_picture_image(prs.slides[14], "图片 3",  os.path.join(IMG_DIR, "p15_chat.png"))
replace_picture_image(prs.slides[14], "图片 27", os.path.join(IMG_DIR, "p15_form.png"))

prs.save(DST)
print(f"saved -> {DST}")
