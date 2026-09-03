# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

path = r"D:\AI\14 - 数据治理\ppt_revision\AI_主数据治理_v3.pptx"
prs = Presentation(path)

for idx in [13]:
    slide = prs.slides[idx-1]
    print(f"### 第 {idx} 页  共 {len(slide.shapes)} 个 shape")
    for sh in slide.shapes:
        stype = sh.shape_type
        info = f"name='{sh.name}' type={stype}"
        try:
            info += f" pos=({sh.left},{sh.top}) size=({sh.width},{sh.height})"
        except Exception:
            pass
        print(info)
        if sh.has_text_frame and sh.text_frame.text.strip():
            print("   text:", repr(sh.text_frame.text))
        if stype == MSO_SHAPE_TYPE.PICTURE:
            print("   >> PICTURE")
        if stype == MSO_SHAPE_TYPE.GROUP:
            print("   >> GROUP, subshapes:")
            for sub in sh.shapes:
                print("      -", sub.name, sub.shape_type,
                      (sub.text_frame.text[:40] if sub.has_text_frame else ""))
