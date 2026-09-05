# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Pt

path = r"D:\AI\14 - 数据治理\ppt_revision\AI_主数据治理_v2.pptx"
prs = Presentation(path)

# target slides: 1-based
targets = {1, 6, 11, 12, 14, 18, 19, 20, 21, 22, 23, 24, 25}

def emu_to_pt(v):
    try:
        return round(v.pt, 1)
    except Exception:
        return v

for idx, slide in enumerate(prs.slides, 1):
    if idx not in targets:
        continue
    print(f"\n{'#'*70}\n# 第 {idx} 页\n{'#'*70}")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        if not tf.text.strip():
            continue
        print(f"\n[Shape] name='{shape.name}'")
        for p_i, para in enumerate(tf.paragraphs):
            runs_info = []
            for r in para.runs:
                f = r.font
                runs_info.append(
                    f"  run(text={r.text!r}, bold={f.bold}, size={emu_to_pt(f.size)}, name={f.name})"
                )
            print(f"  P{p_i} align={para.alignment} level={para.level}")
            for ri in runs_info:
                print(ri)
