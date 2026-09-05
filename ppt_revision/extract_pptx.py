# -*- coding: utf-8 -*-
import sys
from pptx import Presentation
from pptx.util import Emu

path = r"D:\AI\14 - 数据治理\ppt_revision\AI_主数据治理_v2.pptx"
prs = Presentation(path)

print(f"=== 总页数: {len(prs.slides)} ===\n")

for idx, slide in enumerate(prs.slides, 1):
    print(f"\n{'='*60}")
    print(f"### 第 {idx} 页")
    print(f"{'='*60}")
    for shape in slide.shapes:
        # shape name & type
        stype = shape.shape_type
        name = shape.name
        # text frames
        if shape.has_text_frame:
            tf = shape.text_frame
            texts = []
            for para in tf.paragraphs:
                line = "".join(run.text for run in para.runs)
                if not line and para.text:
                    line = para.text
                texts.append(line)
            full = "\n".join(t for t in texts if t is not None)
            if full.strip():
                print(f"\n[文本框 | {name}]")
                print(full)
        # tables
        if shape.has_table:
            tbl = shape.table
            print(f"\n[表格 | {name}] {len(tbl.rows)}行 x {len(tbl.columns)}列")
            for r_i, row in enumerate(tbl.rows):
                cells = []
                for cell in row.cells:
                    cells.append(cell.text.replace("\n", " / "))
                print(" | ".join(cells))
    # notes
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text
        if notes.strip():
            print(f"\n[备注]")
            print(notes)
