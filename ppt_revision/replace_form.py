# -*- coding: utf-8 -*-
"""用 Pillow 绘制的精确表单图替换 P15 图片27，输出 v7"""
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

SRC = r"D:\AI\14 - 数据治理\ppt_revision\AI_主数据治理_v6.pptx"
DST = r"D:\AI\14 - 数据治理\ppt_revision\AI_主数据治理_v7.pptx"
FORM_IMG = r"D:\AI\14 - 数据治理\ppt_revision\new_images\p15_form.png"

prs = Presentation(SRC)
slide = prs.slides[14]  # P15

for sh in slide.shapes:
    if sh.name == "图片 27":
        blip = sh._element.findall('.//' + qn('a:blip'))[0]
        image_part = prs.part.package.get_or_add_image_part(FORM_IMG)
        new_rId = slide.part.relate_to(image_part, RT.IMAGE)
        blip.set(qn('r:embed'), new_rId)
        print("replaced 图片 27 with Pillow-drawn form")
        break

prs.save(DST)
print(f"saved -> {DST}")
