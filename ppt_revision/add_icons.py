# -*- coding: utf-8 -*-
"""
为 P18/P20/P23 卡片页添加图标，调整文字位置：
- P18/P23（高卡片）：图标居中顶部，文字下移并居中
- P20（矮卡片）：图标左侧，文字右移
"""
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import os

SRC = r"D:\AI\14 - 数据治理\ppt_revision\AI_主数据治理_v9.pptx"
DST = r"D:\AI\14 - 数据治理\ppt_revision\AI_主数据治理_v10.pptx"
ICON_DIR = r"D:\AI\14 - 数据治理\ppt_revision\new_images\icons"

prs = Presentation(SRC)

def set_run_font(run, size=None, color=None, bold=None):
    if size: run.font.size = Pt(size)
    if color: run.font.color.rgb = color
    if bold is not None: run.font.bold = bold

def center_text(shape):
    """将文本框所有段落居中"""
    for p in shape.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER

# ============================================================
# P18: 3 张高卡片 — 图标居中顶部
# ============================================================
slide18 = prs.slides[17]
p18_icons = ["icon_multimodal.png", "icon_dedup.png", "icon_formassist.png"]
p18_cards_x = [1.3, 11.9, 22.6]  # 卡片左边距
p18_text_names = ["TextBox 9", "TextBox 12", "TextBox 15"]

for i, (icon_fname, card_x, text_name) in enumerate(zip(p18_icons, p18_cards_x, p18_text_names)):
    # 图标居中顶部
    icon_size = 2.2
    icon_x = card_x + (10.0 - icon_size) / 2
    icon_y = 5.3
    slide18.shapes.add_picture(
        os.path.join(ICON_DIR, icon_fname),
        Cm(icon_x), Cm(icon_y), Cm(icon_size), Cm(icon_size)
    )
    # 文字下移
    for sh in slide18.shapes:
        if sh.name == text_name:
            sh.top = Cm(8.0)
            sh.height = Cm(5.0)
            center_text(sh)
            # 标题字号微调
            break

print("P18 icons added")

# ============================================================
# P20: 5 张矮卡片 — 图标左侧，文字右移
# ============================================================
slide20 = prs.slides[19]

# (icon, card_x, card_y, text_name, text_x_new)
p20_cards = [
    ("icon_llm.png",       1.3,  3.3, "TextBox 8",  3.8),
    ("icon_domain.png",    11.9, 3.3, "TextBox 11", 14.4),
    ("icon_datamodel.png", 22.6, 3.3, "TextBox 14", 25.1),
    ("icon_ontology.png",  1.3,  9.1, "TextBox 17", 3.8),
    ("icon_knowledge.png", 17.2, 9.1, "TextBox 20", 19.7),
]

for icon_fname, card_x, card_y, text_name, text_x_new in p20_cards:
    icon_size = 1.8
    icon_x = card_x + 0.5
    icon_y = card_y + (5.3 - icon_size) / 2
    slide20.shapes.add_picture(
        os.path.join(ICON_DIR, icon_fname),
        Cm(icon_x), Cm(icon_y), Cm(icon_size), Cm(icon_size)
    )
    for sh in slide20.shapes:
        if sh.name == text_name:
            sh.left = Cm(text_x_new)
            sh.width = Cm(8.0)
            break

print("P20 icons added")

# ============================================================
# P23: 3 张高卡片 — 图标居中顶部
# ============================================================
slide23 = prs.slides[22]
p23_icons = ["icon_whitebox.png", "icon_humanloop.png", "icon_shield.png"]
p23_cards_x = [1.3, 11.9, 22.6]
p23_text_names = ["TextBox 8", "TextBox 11", "TextBox 14"]

for icon_fname, card_x, text_name in zip(p23_icons, p23_cards_x, p23_text_names):
    icon_size = 2.2
    icon_x = card_x + (10.0 - icon_size) / 2
    icon_y = 4.4
    slide23.shapes.add_picture(
        os.path.join(ICON_DIR, icon_fname),
        Cm(icon_x), Cm(icon_y), Cm(icon_size), Cm(icon_size)
    )
    for sh in slide23.shapes:
        if sh.name == text_name:
            sh.top = Cm(7.1)
            sh.height = Cm(5.5)
            center_text(sh)
            break

print("P23 icons added")

prs.save(DST)
print(f"saved -> {DST}")
