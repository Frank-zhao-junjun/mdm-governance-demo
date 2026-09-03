# -*- coding: utf-8 -*-
"""修复 P20：改为顶部居中图标布局"""
from pptx import Presentation
from pptx.util import Cm
from pptx.enum.text import PP_ALIGN
import os

SRC = r"D:\AI\14 - 数据治理\ppt_revision\AI_主数据治理_v10.pptx"
DST = r"D:\AI\14 - 数据治理\ppt_revision\AI_主数据治理_v10.pptx"  # 原地覆盖
ICON_DIR = r"D:\AI\14 - 数据治理\ppt_revision\new_images\icons"

prs = Presentation(SRC)
slide20 = prs.slides[19]

# 删除上一版插入的图标图片（通过尺寸和位置识别）
to_remove = []
for sh in slide20.shapes:
    if sh.shape_type == 13:  # PICTURE
        # P20 上一版图标的位置：x=1.8/12.4/23.1/1.8/17.7, y=5.05/5.05/5.05/10.85/10.85
        # 检查是否在卡片区域内的小图片
        if sh.width < Cm(2.5) and sh.height < Cm(2.5):
            to_remove.append(sh)
for sh in to_remove:
    sh._element.getparent().remove(sh._element)
    print(f"removed old icon: {sh.name}")

# 文字框位置恢复并居中
text_info = [
    ("TextBox 8",  1.8,  5.7, 9.1),
    ("TextBox 11", 12.5, 5.7, 9.1),
    ("TextBox 14", 23.1, 5.7, 9.1),
    ("TextBox 17", 1.8,  11.5, 14.5),
    ("TextBox 20", 17.8, 11.5, 14.5),
]
for name, x, y, w in text_info:
    for sh in slide20.shapes:
        if sh.name == name:
            sh.left = Cm(x)
            sh.top = Cm(y)
            sh.width = Cm(w)
            sh.height = Cm(2.5)
            for p in sh.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
            break

# 顶部居中插入图标
p20_cards = [
    ("icon_llm.png",       1.3,  3.3, 10.0),
    ("icon_domain.png",    11.9, 3.3, 10.0),
    ("icon_datamodel.png", 22.6, 3.3, 10.0),
    ("icon_ontology.png",  1.3,  9.1, 15.4),
    ("icon_knowledge.png", 17.2, 9.1, 15.4),
]
for icon_fname, card_x, card_y, card_w in p20_cards:
    icon_size = 1.8
    icon_x = card_x + (card_w - icon_size) / 2
    icon_y = card_y + 0.35
    slide20.shapes.add_picture(
        os.path.join(ICON_DIR, icon_fname),
        Cm(icon_x), Cm(icon_y), Cm(icon_size), Cm(icon_size)
    )
    print(f"added {icon_fname} at ({icon_x:.1f}, {icon_y:.1f})")

prs.save(DST)
print(f"saved -> {DST}")
