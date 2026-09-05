# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

path = r"D:\AI\14 - 数据治理\ppt_revision\AI_主数据治理_v4.pptx"
prs = Presentation(path)

def cm(v):
    try: return round(v/360000, 2)
    except: return v

def fill_info(shape):
    try:
        f = shape.fill
        ft = f.type
        info = f"fill_type={ft}"
        if ft is not None and str(ft) != 'None (0)':
            try:
                if hasattr(f, 'fore_color') and f.fore_color and f.fore_color.type is not None:
                    info += f" color={f.fore_color.rgb}"
            except Exception:
                pass
        return info
    except Exception as e:
        return f"fill_err={e}"

for idx in [15, 16]:
    slide = prs.slides[idx-1]
    print(f"\n{'#'*70}\n# 第 {idx} 页  slide_size={cm(prs.slide_width)}x{cm(prs.slide_height)}cm  shapes={len(slide.shapes)}\n{'#'*70}")
    for sh in slide.shapes:
        st = sh.shape_type
        line = f"name='{sh.name}' type={st} pos=({cm(sh.left)},{cm(sh.top)}) size=({cm(sh.width)}x{cm(sh.height)})"
        print("\n" + line)
        print("  " + fill_info(sh))
        if st == MSO_SHAPE_TYPE.PICTURE:
            print("  >> PICTURE")
        if st == MSO_SHAPE_TYPE.GROUP:
            print("  >> GROUP")
            for sub in sh.shapes:
                print("    -", sub.name, sub.shape_type, cm(sub.left), cm(sub.top), cm(sub.width), cm(sub.height))
        if sh.has_text_frame:
            t = sh.text_frame.text.replace("\n", " / ")
            if t.strip():
                print("  text:", t[:80])
